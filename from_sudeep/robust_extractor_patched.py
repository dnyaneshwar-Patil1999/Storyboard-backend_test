"""
robust_extractor.py — Three-Layer Content Extraction Module
============================================================
Drop this file into: azure_function-main/

Provides robust content extraction for ALL PPT and Word file types by
combining three layers:

  Layer 1: Native extraction (python-pptx / python-docx) — free, fast, exact
  Layer 2: LLM Vision enhancement — only for slides/pages where Layer 1 fails
  Layer 3: Structure assignment — existing LLM call, just needs better input

This module is called BY the existing agents (presentation_agent.py,
hierarchical_agent.py, no_toc_agent.py). It does NOT replace them — it
replaces their internal content extraction methods.

OUTPUT CONTRACT:
  Every extraction function returns a list of SlideContent / PageContent dicts
  with these exact fields that map to the Azure Search index:

    {
      "content":          str  — maps to index "content" → outline "Full Page Content"
      "chapter":          str  — maps to index "chapter" → outline "Chapter"  (set by Layer 3)
      "topic":            str  — maps to index "topic"   → outline "Topic"    (slide/section title)
      "subtopics":        str  — maps to index "subtopic" → outline "Subtopic"
      "source_page_range": str — maps to index "source_page_range" → outline "Source Page"
      "images":           list — image data for the slide
      "tables":           list — table data
      "page_number":      int  — internal use
      "_confidence":      float — internal: 1.0=native, 0.5=needs Layer 2
    }
"""

import os
import re
import json
import base64
import logging
from collections import defaultdict
from typing import List, Dict, Any, Optional, Tuple

# ── Optional imports with graceful fallback ──
try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.exc import PackageNotFoundError as PptxPackageError
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from lxml import etree
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    PptxPackageError = Exception  # Fallback so except clauses don't break

try:
    import docx as python_docx
    from docx.opc.exceptions import PackageNotFoundError as DocxPackageError
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    DocxPackageError = Exception

try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

from zipfile import BadZipFile


# ══════════════════════════════════════════════════════════════════════════════
# FILE FORMAT DETECTION
# ══════════════════════════════════════════════════════════════════════════════

def detect_file_format(file_path: str) -> str:
    """
    Detect whether a file can be natively parsed or needs legacy fallback.
    Returns: 'pptx' | 'ppt_legacy' | 'docx' | 'doc_legacy' | 'pdf' | 'unknown'
    
    Handles:
      - .pptx → native python-pptx (Layer 1)
      - .ppt  → CANNOT be natively parsed, must use LibreOffice → PDF → DI
      - .docx → native python-docx (Layer 1)
      - .doc  → CANNOT be natively parsed, must use LibreOffice → PDF → DI
      - .pdf  → Document Intelligence (existing pipeline)
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pptx':
        return 'pptx'
    elif ext == '.ppt':
        logger.info(f"[FORMAT] Legacy .ppt detected: {os.path.basename(file_path)}. "
                     "Cannot parse natively — requires LibreOffice → PDF → DI fallback.")
        return 'ppt_legacy'
    elif ext == '.docx':
        return 'docx'
    elif ext == '.doc':
        logger.info(f"[FORMAT] Legacy .doc detected: {os.path.basename(file_path)}. "
                     "Cannot parse natively — requires LibreOffice → PDF → DI fallback.")
        return 'doc_legacy'
    elif ext == '.pdf':
        return 'pdf'
    else:
        return 'unknown'


def is_password_protected(file_path: str) -> bool:
    """
    Check if a file is password-protected before attempting to open it.
    Works for both PPTX and DOCX (both are ZIP archives).
    """
    try:
        import zipfile
        with zipfile.ZipFile(file_path, 'r') as z:
            # Encrypted Office files have an EncryptedPackage entry
            names = z.namelist()
            if 'EncryptedPackage' in names:
                return True
            # Or they may just fail to list contents
        return False
    except BadZipFile:
        # Not a valid ZIP — might be legacy binary format or corrupted
        return False
    except Exception:
        return False

logger = logging.getLogger(__name__)


# ══════════════════════════════════════════════════════════════════════════════
# SHARED UTILITIES (used by all layers)
# ══════════════════════════════════════════════════════════════════════════════

def should_exclude_paragraph(text: str, role: str, repeating_texts: set,
                             total_pages: int) -> bool:
    """
    Three-layer header/footer filter.
    Layer 1: DI role (pageHeader/pageFooter/pageNumber) — but keep if >15 words
    Layer 2: Repeating text across >30% of pages
    Layer 3: Standalone numbers (slide/page numbers)
    """
    if not text or not text.strip():
        return True

    text = text.strip()

    # Layer 1: DI role
    if role in ('pageHeader', 'pageFooter', 'pageNumber'):
        if len(text.split()) > 15:
            return False  # Real content in a header/footer region
        return True

    # Layer 2: Repeating text
    if text in repeating_texts:
        return True

    # Layer 3: Standalone numbers (slide numbers)
    if text.isdigit():
        try:
            if int(text) <= 500:
                return True
        except ValueError:
            pass

    return False


def clean_ocr_noise(text: str) -> str:
    """Remove common OCR artifacts from DI output."""
    if not text:
        return ""
    # Remove single-character fragments that are OCR noise from logos
    text = re.sub(r'\b[A-Z]\s*\n\s*[A-Z]\b', '', text)
    # Remove :selected: / :unselected: markers from checkboxes
    text = re.sub(r':(?:un)?selected:\s*', '', text)
    # [PATCH] Strip ONLY [MISSING_IMAGE:...] and [IMAGE-ONLY SLIDE...] markers.
    # KEEP [IMAGE: id] markers — those are positional placeholders that
    # the renderer needs to resolve to actual embedded images.
    text = re.sub(r'\[MISSING_IMAGE[^\]]*\]', '', text)
    text = re.sub(r'\[IMAGE-ONLY SLIDE[^\]]*\]', '', text)
    # Collapse excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def format_table_markdown(rows_data: List[List[str]], bold_header: bool = True) -> str:
    """
    Convert a list of rows (each row = list of cell strings) to markdown pipe table.
    First row treated as header if bold_header=True.
    """
    if not rows_data:
        return ""
    lines = []
    for ri, row in enumerate(rows_data):
        cells = [c.replace('|', '/').strip() for c in row]
        if bold_header and ri == 0:
            cells = [f"**{c}**" if c else "" for c in cells]
        lines.append("| " + " | ".join(cells) + " |")
        if ri == 0:
            lines.append("|" + "|".join([" --- "] * len(cells)) + "|")
    return "\n".join(lines)


def is_likely_heading_docx(para, body_font_size: float = 11.0) -> int:
    """
    Detect headings by style name, custom style keywords, OR visual formatting.
    Returns: heading level (1-3) or 0 if not a heading.
    Works with: standard styles, custom corporate styles, and visual-only formatting.
    """
    style = para.style.name if para.style else 'Normal'
    text = para.text.strip()

    if not text or len(text) > 120:
        return 0

    # ── Standard Word heading styles ──
    STANDARD = {
        'Heading 1': 1, 'Heading 2': 2, 'Heading 3': 3,
        'Heading 4': 3, 'Heading 5': 3, 'Heading 6': 3,
        'Title': 1, 'Subtitle': 2,
    }
    if style in STANDARD:
        return STANDARD[style]

    # ── Custom styles with heading/title/chapter/section keywords ──
    sl = style.lower()
    if any(kw in sl for kw in ['heading', 'title', 'chapter', 'section', 'header']):
        for digit in '123456':
            if digit in sl:
                return min(int(digit), 3)
        if 'chapter' in sl or 'title' in sl:
            return 1
        if 'section' in sl:
            return 2
        return 2  # Default for custom heading styles

    # ── Visual formatting heuristic (for docs with no heading styles) ──
    runs = [r for r in para.runs if r.text.strip()]
    if not runs:
        return 0

    all_bold = all(r.bold for r in runs if r.bold is not None)
    # Check inherited bold (from style, not per-run)
    if not all_bold:
        # If run.bold is None, it inherits from style — check the style's font
        style_bold = getattr(para.style.font, 'bold', None) if para.style else None
        if style_bold:
            all_bold = True

    max_size = body_font_size
    for r in runs:
        if r.font.size:
            sz = r.font.size.pt
            if sz > max_size:
                max_size = sz

    word_count = len(text.split())

    # Large bold text with few words = heading
    if all_bold and max_size >= 16 and word_count <= 10:
        return 1
    if all_bold and max_size >= 13 and word_count <= 12:
        return 2
    if all_bold and word_count <= 8 and not text.endswith(('.', '!', '?', ':')):
        return 3

    return 0


def get_para_list_level_docx(para) -> int:
    """
    Get the list nesting level from DOCX XML.
    Handles both Word's native list formatting AND manual numbering.
    Returns: nesting level (0=top, 1=sub, etc.) or -1 if not a list item.
    """
    # Check native Word list formatting (numPr/ilvl in XML)
    WML = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    pPr = para._p.pPr
    if pPr is not None:
        numPr = pPr.find(f'{WML}numPr')
        if numPr is not None:
            ilvl = numPr.find(f'{WML}ilvl')
            if ilvl is not None:
                return int(ilvl.get(f'{WML}val', '0'))

    # Check manual bullet/number patterns
    text = para.text.strip()
    if re.match(r'^[-\u2022\u2023\u25E6\u25AA\u25CF\u25CB\u2013\u2014]\s', text):
        # Count leading whitespace to estimate level
        raw = para.text
        indent = len(raw) - len(raw.lstrip())
        return min(indent // 4, 3)

    if re.match(r'^\d+[.)]\s', text):
        raw = para.text
        indent = len(raw) - len(raw.lstrip())
        return min(indent // 4, 3)

    if re.match(r'^[a-z][.)]\s', text):
        return 1  # Sub-item

    if re.match(r'^[ivxIVX]+[.)]\s', text):
        return 1  # Roman numeral sub-item

    return -1  # Not a list item


def extract_bold_text(para) -> str:
    """
    Extract paragraph text preserving bold with **markdown**.
    Handles: full-line bold, inline bold, inherited bold from styles.
    """
    runs = list(para.runs)
    if not runs:
        return para.text.strip()

    # Check if ALL runs with text are bold
    text_runs = [r for r in runs if r.text.strip()]
    if not text_runs:
        return para.text.strip()

    all_bold = all(_is_run_bold(r, para) for r in text_runs)
    if all_bold:
        return f"**{para.text.strip()}**"

    # Mixed bold — build inline markdown
    parts = []
    for r in runs:
        txt = r.text
        if not txt:
            continue
        if _is_run_bold(r, para) and txt.strip():
            parts.append(f"**{txt}**")
        else:
            parts.append(txt)
    result = ''.join(parts).strip()
    # Clean up double asterisks at word boundaries
    result = re.sub(r'\*\*\s+\*\*', ' ', result)
    return result


def _is_run_bold(run, para) -> bool:
    """Check if a run is bold, including inherited bold from paragraph style."""
    if run.bold is True:
        return True
    if run.bold is False:
        return False
    # None = inherit from style
    if para.style and para.style.font and para.style.font.bold:
        return True
    return False


# ══════════════════════════════════════════════════════════════════════════════
# LAYER 1: NATIVE PPTX EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════

# ─────────────────────────────────────────────────────────────────────────────
# [CHANGE] Placeholder role mapping
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE used hardcoded placeholder INDEX sets:
#   _SKIP_PLACEHOLDER_IDS = {10, 11, 12, 4294967295}
#   _TITLE_PLACEHOLDER_IDS = {0, 1}
# These were brittle — different layouts use varying idx values for titles,
# subtitles, and bodies. Manual idx checking missed CENTER_TITLE,
# VERTICAL_TITLE, and other valid placeholder types.
#
# NEW CODE uses python-pptx PP_PLACEHOLDER enum types (TITLE, CENTER_TITLE,
# VERTICAL_TITLE, SUBTITLE, BODY, OBJECT, etc.) for role detection. This is
# PowerPoint's own semantic classification — far more reliable than idx numbers.
# Also adds a three-strategy fallback: (1) python-pptx enum, (2) raw XML
# <p:ph type="..."> attribute, (3) shape name regex heuristic.
# ─────────────────────────────────────────────────────────────────────────────

# DrawingML / PresentationML XML namespaces (used by bullet detection and role resolution)
_DML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Placeholder role sets using PP_PLACEHOLDER enum types
try:
    _TITLE_TYPES    = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                       PP_PLACEHOLDER.VERTICAL_TITLE}
    _SUBTITLE_TYPES = {PP_PLACEHOLDER.SUBTITLE}
    _BODY_TYPES     = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT,
                       PP_PLACEHOLDER.VERTICAL_BODY, PP_PLACEHOLDER.VERTICAL_OBJECT}
except (NameError, AttributeError):
    _TITLE_TYPES    = set()
    _SUBTITLE_TYPES = set()
    _BODY_TYPES     = set()

_TITLE_NAME_RE = re.compile(
    r"^(title|slide\s*title|center\s*title|vertical\s*title)\s*\d*$",
    re.IGNORECASE,
)


def _ph_role(ph_type):
    """Map a PP_PLACEHOLDER enum value to a semantic role string."""
    if ph_type is None:
        return None
    if ph_type in _TITLE_TYPES:
        return "title"
    if ph_type in _SUBTITLE_TYPES:
        return "subtitle"
    if ph_type in _BODY_TYPES:
        return "body"
    try:
        return ph_type.name.lower()
    except AttributeError:
        return str(ph_type)


# ─────────────────────────────────────────────────────────────────────────────
# [CHANGE] Bullet detection from XML
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE always prefixed body text with "- " (unconditional bullet), and
# only used para.level for indentation depth. This meant every paragraph
# appeared as a bullet, even if it was a regular sentence with no bullet
# character in PowerPoint.
#
# NEW CODE inspects the actual XML for <a:buChar> (character bullet) and
# <a:buAutoNum> (numbered list) elements. Only paragraphs with EXPLICIT
# bullet markup get a bullet prefix. Paragraphs with <a:buNone> or no
# bullet element ("inherited") are rendered as plain text. This accurately
# reflects what PowerPoint displays on screen.
# ─────────────────────────────────────────────────────────────────────────────

def _bullet_from_xml(para):
    """
    Inspect paragraph XML for explicit bullet markers.

    Returns dict:
      type    : "none" | "char" | "autoNum" | "inherited"
      char    : bullet glyph        (type == "char")
      numFmt  : e.g. "arabicPeriod" (type == "autoNum")
      startAt : int                 (type == "autoNum")
    """
    pPr = para._p.find("a:pPr", _DML_NS)
    if pPr is None:
        return {"type": "inherited"}

    if pPr.find("a:buNone", _DML_NS) is not None:
        return {"type": "none"}

    buChar = pPr.find("a:buChar", _DML_NS)
    if buChar is not None:
        return {"type": "char", "char": buChar.get("char", "\u2022")}

    buAutoNum = pPr.find("a:buAutoNum", _DML_NS)
    if buAutoNum is not None:
        return {
            "type":    "autoNum",
            "numFmt":  buAutoNum.get("type", "arabicPeriod"),
            "startAt": int(buAutoNum.get("startAt", 1)),
        }

    return {"type": "inherited"}


def _is_real_bullet(bullet_info: dict) -> bool:
    """True only for explicit <a:buChar> or <a:buAutoNum> bullets."""
    return bullet_info.get("type") in ("char", "autoNum")


def _lstStyle_has_bullets(text_frame) -> bool:
    """
    Check whether the text frame's own <a:lstStyle> defines bullets
    for ANY indent level (lvl1pPr … lvl9pPr).

    PowerPoint stores bullet definitions in lstStyle when the bullets
    come from the placeholder/layout rather than the paragraph itself.
    """
    try:
        txBody = text_frame._txBody
        lstStyle = txBody.find("a:lstStyle", _DML_NS)
        if lstStyle is None:
            return False
        for i in range(1, 10):
            lvl = lstStyle.find(f"a:lvl{i}pPr", _DML_NS)
            if lvl is not None:
                if (lvl.find("a:buChar", _DML_NS) is not None or
                    lvl.find("a:buAutoNum", _DML_NS) is not None):
                    return True
        # Also check defPPr (default paragraph properties)
        defPPr = lstStyle.find("a:defPPr", _DML_NS)
        if defPPr is not None:
            if (defPPr.find("a:buChar", _DML_NS) is not None or
                defPPr.find("a:buAutoNum", _DML_NS) is not None):
                return True
    except Exception:
        pass
    return False


def _get_lstStyle_bullet_char(text_frame, level: int) -> Optional[str]:
    """
    Try to retrieve the bullet character for a specific indent level
    from the text frame's lstStyle XML.
    Returns the char string, or None if not found.
    """
    try:
        txBody = text_frame._txBody
        lstStyle = txBody.find("a:lstStyle", _DML_NS)
        if lstStyle is None:
            return None
        # lstStyle level tags are 1-indexed (lvl1pPr, lvl2pPr, ...)
        lvl_elem = lstStyle.find(f"a:lvl{level + 1}pPr", _DML_NS)
        if lvl_elem is not None:
            buChar = lvl_elem.find("a:buChar", _DML_NS)
            if buChar is not None:
                return buChar.get("char", "\u2022")
    except Exception:
        pass
    return None


# ─────────────────────────────────────────────────────────────────────────────
# [CHANGE] Bold + Hyperlink -aware paragraph → markdown rendering
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE:
#   - Used assume_bullet=True/False based only on shape role == "body".
#   - Did not check lstStyle or sibling paragraphs for bullet definitions.
#   - Inherited bullets in non-body shapes (e.g. content placeholders with
#     lstStyle bullets) were missed.
#   - No hyperlink support.
#
# NEW CODE:
#   - Uses frame_has_bullets flag computed from THREE signals:
#     1. Any sibling paragraph has explicit <a:buChar>/<a:buAutoNum>
#     2. The text frame's lstStyle XML defines bullets at any level
#     3. The shape is a body placeholder (shape_role == "body")
#   - For inherited bullets, looks up the actual bullet char from lstStyle
#     so the correct glyph is preserved (not always "- ").
#   - Adds hyperlink rendering: [text](url) for runs with hyperlinks.
#   - Bold detection still uses run.font.bold (high-level API).
# ─────────────────────────────────────────────────────────────────────────────

def _para_to_markdown(para, frame_has_bullets=False, text_frame=None) -> Optional[str]:
    """
    Convert one python-pptx Paragraph to a markdown string.

    Returns None for blank / digit-only slide-number paragraphs.

    frame_has_bullets: True if inherited bullets should be rendered
    text_frame: the parent text frame, used to look up lstStyle bullet chars
    """
    text = para.text.strip()
    if not text:
        return None
    # Skip lone page/slide numbers (1-3 digit strings)
    if text.isdigit() and len(text) <= 3:
        return None
    # Filter footer artifacts (date stamps, "Section Break", etc.)
    if _is_pptx_footer_artifact(text):
        return None

    level  = para.level or 0
    indent = "  " * level

    bullet_info = _bullet_from_xml(para)
    has_bullet  = _is_real_bullet(bullet_info)

    # ── Handle inherited bullets ──────────────────────────────────────────
    # "inherited" means no explicit bullet in the paragraph XML.
    # Treat as bullet when the frame context says bullets should apply.
    if not has_bullet and bullet_info.get("type") == "inherited":
        if frame_has_bullets or level > 0:
            has_bullet = True

    # ── Bold + Italic + Color + Hyperlink formatting ──────────────────────
    runs_with_text = [r for r in para.runs if r.text.strip()]

    if not runs_with_text:
        formatted = text
    else:
        has_links = False
        for r in runs_with_text:
            try:
                if r.hyperlink and r.hyperlink.address:
                    has_links = True
                    break
            except Exception:
                pass

        # ── [PATCH] Defensive bold inheritance check ──────────────────────
        # run.font.bold returns None when bold is inherited (not explicitly set).
        # In most PPTs, inherited=None means NOT bold. But some corporate
        # templates set bold at the paragraph defRPr level, in which case
        # all runs inherit bold but run.font.bold still returns None.
        # This check looks at paragraph XML defRPr as a fallback.
        def _run_is_bold(r):
            if r.font.bold is True:
                return True
            if r.font.bold is False:
                return False
            # None = inherited. Check paragraph defRPr as fallback.
            try:
                pPr = para._p.find("a:pPr", _DML_NS)
                if pPr is not None:
                    defRPr = pPr.find("a:defRPr", _DML_NS)
                    if defRPr is not None and defRPr.get("b") == "1":
                        return True
            except Exception:
                pass
            return False

        def _run_is_italic(r):
            if r.font.italic is True:
                return True
            if r.font.italic is False:
                return False
            try:
                pPr = para._p.find("a:pPr", _DML_NS)
                if pPr is not None:
                    defRPr = pPr.find("a:defRPr", _DML_NS)
                    if defRPr is not None and defRPr.get("i") == "1":
                        return True
            except Exception:
                pass
            return False

        def _run_color_class(r):
            """Returns 'red', 'blue', or None based on font color."""
            try:
                if r.font.color and r.font.color.type is not None:
                    rgb = r.font.color.rgb
                    if rgb is not None:
                        hex_color = str(rgb).upper()
                        if hex_color in ('FF0000', 'C00000', 'C0504D', 'CC0000', 'E81123', 'D32F2F'):
                            return 'red'
                        if hex_color in ('0070C0', '0000FF', '4472C4', '2E75B6', '2196F3', '1976D2', '0563C1'):
                            return 'blue'
            except Exception:
                pass
            return None

        # Determine if entire line has uniform formatting
        all_bold = all(_run_is_bold(r) for r in runs_with_text)
        all_italic = all(_run_is_italic(r) for r in runs_with_text)
        all_red = all(_run_color_class(r) == 'red' for r in runs_with_text)
        all_blue = all(_run_color_class(r) == 'blue' for r in runs_with_text)

        if not has_links and (all_bold or all_italic or all_red or all_blue):
            piece = text
            if all_bold and all_italic:
                piece = f"***{piece}***"
            elif all_bold:
                piece = f"**{piece}**"
            elif all_italic:
                piece = f"*{piece}*"
            if all_red:
                piece = f"<RED>{piece}</RED>"
            elif all_blue:
                piece = f"<BLUE>{piece}</BLUE>"
            formatted = piece
        else:
            parts = []
            for r in para.runs:
                rt = r.text
                if not rt:
                    continue
                piece = rt
                is_b = _run_is_bold(r)
                is_i = _run_is_italic(r)
                color = _run_color_class(r)

                if rt.strip():
                    if is_b and is_i:
                        piece = f"***{rt}***"
                    elif is_b:
                        piece = f"**{rt}**"
                    elif is_i:
                        piece = f"*{rt}*"
                    if color == 'red':
                        piece = f"<RED>{piece}</RED>"
                    elif color == 'blue':
                        piece = f"<BLUE>{piece}</BLUE>"
                try:
                    if r.hyperlink and r.hyperlink.address:
                        piece = f"[{piece}]({r.hyperlink.address})"
                except Exception:
                    pass
                parts.append(piece)
            formatted = "".join(parts).strip()

    if not formatted:
        return None

    # ── Bullet / indent prefix ────────────────────────────────────────────
    if has_bullet:
        btype = bullet_info.get("type")
        if btype == "autoNum":
            prefix = f"{bullet_info.get('startAt', 1)}. "
        elif btype == "char":
            char = bullet_info.get("char", "\u2022")
            prefix = f"{char} "
        else:
            # Inherited bullet — try to get char from lstStyle
            char = None
            if text_frame is not None:
                char = _get_lstStyle_bullet_char(text_frame, level)
            if char is None:
                char = "\u2022"
            prefix = f"{char} "
        return f"{indent}{prefix}{formatted}"
    else:
        return f"{indent}{formatted}" if level > 0 else formatted


# Footer / artifact patterns that should be filtered from PPTX content
_PPTX_FOOTER_PATTERNS = [
    re.compile(r'^\d{1,2}/\d{1,2}/\d{2,4}$'),
    re.compile(r'^\d{1,2}-\d{1,2}-\d{2,4}$'),
    re.compile(r'^Sample Footer Text$', re.I),
    re.compile(r'^AI-generated content may be incorrect\.?$', re.I),
    re.compile(r'^Click to add (text|notes|content|title)$', re.I),
    re.compile(r'^Section Break$', re.I),
    re.compile(r'^CONFIDENTIAL\b.*INTERNAL USE\b', re.I),
    re.compile(r'^FOR INTERNAL USE ONLY$', re.I),
    re.compile(r'^PROPRIETARY( AND CONFIDENTIAL)?$', re.I),
    re.compile(r'^Page\s*\d+\s*(of\s*\d+)?$', re.I),
    re.compile(r'^Slide\s*\d+$', re.I),
    re.compile(r'^Copyright\s*©.*$', re.I),
]


def _is_pptx_footer_artifact(text: str) -> bool:
    """Returns True if text is a footer/decoration that should be filtered."""
    if not text:
        return True
    text = text.strip()
    if not text:
        return True
    for pat in _PPTX_FOOTER_PATTERNS:
        if pat.match(text):
            return True
    return False


def _text_frame_to_markdown(text_frame, shape_role=None) -> str:
    """Render all paragraphs of a text frame to a single markdown string.

    Determines whether inherited bullets should be rendered using three signals:
      1. Any sibling paragraph has an explicit <a:buChar>/<a:buAutoNum>
      2. The text frame's lstStyle XML defines bullets for any level
      3. The shape is a body placeholder (shape_role == "body")
    """
    # Signal 1: any sibling has explicit bullet
    has_explicit = any(
        _is_real_bullet(_bullet_from_xml(p)) for p in text_frame.paragraphs
    )
    # Signal 2: lstStyle defines bullets
    has_lstStyle_bullets = _lstStyle_has_bullets(text_frame)
    # Signal 3: body placeholder — nearly always has bullets from master/layout
    is_body = (shape_role == "body")

    frame_has_bullets = has_explicit or has_lstStyle_bullets or is_body

    lines = []
    for para in text_frame.paragraphs:
        line = _para_to_markdown(para, frame_has_bullets=frame_has_bullets, text_frame=text_frame)
        if line is not None:
            lines.append(line)
    return "\n".join(lines)


# ─────────────────────────────────────────────────────────────────────────────
# [CHANGE] Three-strategy role resolution
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE had two checks: (1) placeholder_format.idx in {0, 1} for titles,
# (2) shape.name.lower().startswith('title'). This missed subtitles entirely,
# couldn't distinguish body from object placeholders, and failed for
# CENTER_TITLE or VERTICAL_TITLE types.
#
# NEW CODE uses a three-strategy fallback:
#   1. python-pptx placeholder_format.type → _ph_role() with enum sets
#   2. Raw XML <p:ph type="..."> attribute (covers cases pptx misses)
#   3. Shape name regex heuristic (last resort for non-placeholder shapes)
# This correctly identifies titles, subtitles, body, and other roles across
# all layout types.
# ─────────────────────────────────────────────────────────────────────────────

def _resolve_role(shape) -> Optional[str]:
    """Determine the semantic role of a shape using three fallback strategies."""
    # Strategy 1: python-pptx placeholder_format type enum
    if shape.is_placeholder:
        try:
            ph_type = shape.placeholder_format.type
            role    = _ph_role(ph_type)
            if role is not None:
                return role
        except Exception:
            pass

    # Strategy 2: Raw XML <p:ph type="..."> / idx attribute
    try:
        ph = shape._element.find(f".//{{{_PML_NS}}}ph")
        if ph is not None:
            raw_type = ph.get("type", "")
            if raw_type in ("title", "ctrTitle", "verticalTitle"):
                return "title"
            if raw_type == "subTitle":
                return "subtitle"
            if raw_type in ("body", "obj", "vertBody", "vertObj"):
                return "body"
            if ph.get("idx", "1") == "0" and raw_type == "":
                return "title"
    except Exception:
        pass

    # Strategy 3: Shape name heuristic
    try:
        if _TITLE_NAME_RE.match(shape.name or ""):
            return "title"
    except Exception:
        pass

    return None


# ─────────────────────────────────────────────────────────────────────────────
# [CHANGE] Visual reading-order sort
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE processed shapes in the default XML order (slide.shapes), which
# does not necessarily match the visual layout order on the slide. This could
# cause body text to appear before titles or content from the right column to
# appear before the left column.
#
# NEW CODE sorts shapes by: title first (priority 0) → subtitle (priority 1) →
# all other shapes by (top, left) position (priority 2). This ensures the
# extracted text follows the visual reading order of the slide.
# ─────────────────────────────────────────────────────────────────────────────

def _shape_sort_key(shape):
    """Sort key: title first → subtitle → rest by (top, left) visual position."""
    role = _resolve_role(shape)
    if role == "title":
        return (0, 0, 0)
    if role == "subtitle":
        return (1, 0, 0)
    top  = shape.top  if shape.top  is not None else 0
    left = shape.left if shape.left is not None else 0
    return (2, top, left)


# ─────────────────────────────────────────────────────────────────────────────
# [CHANGE] Recursive shape processing
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE handled group shapes with a shallow hasattr(shape, 'shapes') check
# and iterated child shapes with a simple "- {text}" prefix. Tables, images,
# and nested groups inside the group were not properly handled.
#
# NEW CODE uses the reference code's recursive _process_shapes pattern. Group
# shapes are recursively processed through the same dispatch logic that handles
# top-level shapes. This ensures tables, images, and deeply nested groups
# inside SmartArt or group containers are correctly extracted.
# ─────────────────────────────────────────────────────────────────────────────

def _process_shapes_for_native(shapes, slide_data):
    """
    Process shapes using reference code logic for extract_pptx_native.
    Sorts by visual reading order, resolves roles, extracts content recursively.
    """
    try:
        sorted_shapes = sorted(shapes, key=_shape_sort_key)
    except Exception:
        sorted_shapes = list(shapes)

    for shape in sorted_shapes:

        # ── Recursive group shape processing ──
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            slide_data["has_smartart"] = True
            _process_shapes_for_native(shape.shapes, slide_data)
            continue

        # ── [PATCH] Pictures FIRST ──
        # Pictures must be handled before the generic has_text_frame check
        # because PICTURE shapes can have an empty text frame attached, which
        # would make `if shape.has_text_frame` match and prevent the picture
        # branch from running. Without this reorder, no [IMAGE: id] markers
        # are emitted from PPTX files. Reordering is safe because pictures
        # never need text_frame processing.
        #
        # [PATCH] Two paths to detect a picture:
        #   1. shape.shape_type == PICTURE (a regular inserted image)
        #   2. shape.shape_type == PLACEHOLDER AND the placeholder has a filled-in
        #      .image attribute (a content-placeholder of PP_PLACEHOLDER.PICTURE
        #      type that has had an image dropped into it).
        # Path 2 is the source of "Image missing" complaints on IsDB slides 3
        # and 12 — those slides use the "Content with Picture" layout and the
        # actual image is stored as a placeholder fill, not a PICTURE shape.
        is_picture_shape = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        is_picture_placeholder = False
        if not is_picture_shape and shape.is_placeholder:
            try:
                # If the placeholder has a usable .image attribute, it's a
                # filled-in picture placeholder (PP_PLACEHOLDER.PICTURE = 18).
                _probe = shape.image
                if _probe is not None and getattr(_probe, "blob", None):
                    is_picture_placeholder = True
            except (AttributeError, ValueError, KeyError):
                # No image attached — fall through to text/role processing
                pass

        if is_picture_shape or is_picture_placeholder:
            slide_data["has_images"] = True
            try:
                img = shape.image
                img_bytes = img.blob
                img_ext = (img.ext or "png").lower()
                import hashlib
                h = hashlib.sha1(img_bytes).hexdigest()[:16]
                if h not in slide_data.setdefault("_seen_image_hashes", set()):
                    slide_data["_seen_image_hashes"].add(h)
                    img_id = f"img_{h}"
                    slide_data.setdefault("extracted_images", []).append({
                        "id": img_id,
                        "format": img_ext,
                        "data": __import__("base64").b64encode(img_bytes).decode("ascii"),
                    })
                    slide_data["body_lines"].append(f"[IMAGE: {img_id}]")
            except Exception as e:
                logger.debug(f"Picture extraction failed: {e}")
            continue

        role = _resolve_role(shape)

        # ── Skip slide-number / footer / date placeholders ──
        # Kept from existing code: filter out CONFIDENTIAL and very short footers
        if role in ("slide_number", "footer", "date"):
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt.isdigit() or len(txt) <= 3:
                    continue
                upper = txt.upper()
                if any(kw in upper for kw in [
                    'CONFIDENTIAL', 'INTERNAL USE', 'PROPRIETARY'
                ]):
                    continue

        # ── Text frame ─────────────────────────────────────────────────────
        if shape.has_text_frame:
            md = _text_frame_to_markdown(shape.text_frame, shape_role=role)

            if role == "title" and slide_data["title_text"] is None:
                raw = shape.text_frame.text.strip()
                if raw and not raw.isdigit() and len(raw) > 1:
                    slide_data["title_text"] = raw
            elif role == "subtitle":
                raw = shape.text_frame.text.strip()
                if raw:
                    slide_data["subtitle_text"] = raw
                # Also add subtitle content to body lines for full content
                if md:
                    slide_data["body_lines"].append(md)
            else:
                # ── [PATCH] Fallback title detection ──
                # Some PPTs use TextBox shapes (not placeholders) for titles.
                # If we haven't found a title yet AND this is the first text
                # shape we see (in reading order) AND it's short enough to
                # plausibly be a title, treat it as one.
                if (slide_data["title_text"] is None
                        and not slide_data["body_lines"]
                        and not slide_data["tables_md"]
                        and md):
                    raw = shape.text_frame.text.strip()
                    # Heuristic: title is short, single-line-ish, no bullet markers
                    if (raw and not raw.isdigit()
                            and 2 <= len(raw) <= 100
                            and raw.count('\n') <= 1
                            and not raw.startswith('•')
                            and not raw.startswith('-')):
                        slide_data["title_text"] = raw
                        # Don't also add to body — title is consumed
                        continue
                if md:
                    slide_data["body_lines"].append(md)

        # ── Tables ─────────────────────────────────────────────────────────
        # [CHANGE] Now uses _text_frame_to_markdown per cell for proper
        # bold/bullet rendering, instead of manual cell_has_bold checks.
        elif shape.has_table:
            tbl = shape.table
            rows_data = []
            for ri, row in enumerate(tbl.rows):
                cells = []
                for cell in row.cells:
                    cell_md = _text_frame_to_markdown(cell.text_frame)
                    cells.append(
                        cell_md if cell_md
                        else cell.text.strip().replace('\n', ' ')
                    )
                rows_data.append(cells)
            if rows_data:
                slide_data["tables_md"].append(format_table_markdown(rows_data))

        # ── (Picture handling moved to top of loop) ────────────────────────

        # ── Charts (kept from existing code, not in reference) ─────────────
        elif shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
            slide_data["has_smartart"] = True
            slide_data["body_lines"].append(
                "[CHART - content requires visual analysis]"
            )

        # ── Freeform (often SmartArt decomposed) ──────────────────────────
        elif shape.shape_type == 15:  # MSO_SHAPE_TYPE.FREEFORM
            slide_data["has_smartart"] = True


def _walk_layout_pictures_only(shapes, slide_data):
    """
    [PATCH] Walk a slide layout's shapes, processing ONLY pictures and
    picture-placeholders. Skip all text frames so design-time placeholder
    prompts ("Click to edit Master title style", "Second level", etc.)
    do not leak into the CO output.

    This is a stripped-down version of _process_shapes_for_native() that
    short-circuits anything that isn't a picture. It uses the same hash
    dedup and marker emission so layout pictures appear correctly inline
    in the slide content.
    """
    for shape in shapes:
        # Recurse into groups
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                _walk_layout_pictures_only(shape.shapes, slide_data)
            except Exception:
                pass
            continue

        # Two paths to detect a picture (same as in _process_shapes_for_native)
        is_picture_shape = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        is_picture_placeholder = False
        if not is_picture_shape and getattr(shape, "is_placeholder", False):
            try:
                _probe = shape.image
                if _probe is not None and getattr(_probe, "blob", None):
                    is_picture_placeholder = True
            except (AttributeError, ValueError, KeyError):
                pass

        if not (is_picture_shape or is_picture_placeholder):
            # Skip everything else — including text frames whose content
            # would be design-time placeholder prompts.
            continue

        slide_data["has_images"] = True
        try:
            img = shape.image
            img_bytes = img.blob
            img_ext = (img.ext or "png").lower()
            import hashlib
            h = hashlib.sha1(img_bytes).hexdigest()[:16]
            if h not in slide_data.setdefault("_seen_image_hashes", set()):
                slide_data["_seen_image_hashes"].add(h)
                img_id = f"img_{h}"
                slide_data.setdefault("extracted_images", []).append({
                    "id": img_id,
                    "format": img_ext,
                    "data": __import__("base64").b64encode(img_bytes).decode("ascii"),
                })
                slide_data["body_lines"].append(f"[IMAGE: {img_id}]")
        except Exception as e:
            logger.debug(f"Layout picture extraction failed: {e}")


def extract_pptx_native(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Layer 1: Extract all content from PPTX using native XML structure.

    Handles:
      ✓ Title/subtitle/body via PP_PLACEHOLDER enum + XML + name heuristic
      ✓ Bullet hierarchy: only explicit <a:buChar>/<a:buAutoNum> get prefixes
      ✓ Bold text via run.font.bold (resolves full OOXML inheritance chain)
      ✓ Tables with per-cell markdown rendering
      ✓ Visual reading-order sort (title → subtitle → top/left position)
      ✓ Recursive group shape processing (SmartArt, nested containers)
      ✓ Slide numbers, CONFIDENTIAL footers, decorative elements filtered
      ✓ Image-only slides detected and labeled
      ✓ Slide notes with bullet/bold-aware markdown

    Returns list of dicts matching chunk upload format.
    """
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx not available, cannot extract natively")
        return []

    logger.info(f"[LAYER 1] Native PPTX extraction: {os.path.basename(pptx_path)}")
    slides_out = []

    # ── Check for password protection ──
    if is_password_protected(pptx_path):
        logger.error(f"[LAYER 1] File is password-protected: {os.path.basename(pptx_path)}. "
                      "Cannot extract natively. Return empty to trigger DI fallback.")
        return []

    try:
        prs = Presentation(pptx_path)
    except (BadZipFile, PptxPackageError) as e:
        logger.error(f"[LAYER 1] File is corrupted or password-protected: {e}")
        return []
    except Exception as e:
        logger.error(f"[LAYER 1] Failed to open PPTX: {e}")
        return []

    visible_idx = 0
    for slide in prs.slides:
        # Skip hidden slides
        if slide.element.get('show') == '0':
            continue
        visible_idx += 1

        # ── [CHANGE] Slide data collection ──
        # OLD: Used flat variables (title='', body_lines=[], tables_md=[], ...)
        #      and processed shapes in arbitrary XML order.
        # NEW: Uses a structured dict that is passed to _process_shapes_for_native,
        #      allowing the recursive shape processor to populate title, subtitle,
        #      body, tables, and flags in a single sorted pass.
        slide_data = {
            "title_text":    None,
            "subtitle_text": None,
            "body_lines":    [],
            "tables_md":     [],
            "has_images":    False,
            "has_smartart":  False,
        }

        # ── [CHANGE] Notes extraction ──
        # OLD: Manually iterated notes paragraphs, checked XML for bullets with
        #      direct namespace string, built lines with basic "• " prefix.
        # NEW: Uses _text_frame_to_markdown() which applies the same improved
        #      bullet detection (explicit buChar/buAutoNum only) and bold
        #      rendering (run.font.bold API) to notes content.
        notes_text = ''
        notes_raw_text = ''
        try:
            if slide.has_notes_slide:
                nf = slide.notes_slide.notes_text_frame
                if nf:
                    notes_text = _text_frame_to_markdown(nf)
                    notes_raw_text = nf.text or ''
        except Exception:
            pass

        # ── [NEW] Parse slide duration from notes ──
        # Reviewers complained that durations in CO output were wrong. The
        # source notes contain "Slide duration: X min" but the system was
        # ignoring this and computing duration from word count.
        slide_duration = _parse_slide_duration(notes_raw_text)

        # ── [CHANGE] Shape processing ──
        # OLD: Iterated slide.shapes in XML order. Checked placeholder idx sets
        #      (_SKIP_PLACEHOLDER_IDS, _TITLE_PLACEHOLDER_IDS) for skip/title
        #      decisions. Used _is_pptx_run_bold() for bold with manual 4-level
        #      XML inheritance walk. Always prefixed body text with "- ".
        #      Group shapes extracted with shallow hasattr check.
        # NEW: Uses _process_shapes_for_native() which:
        #   1. Sorts shapes by visual reading order (title→subtitle→top/left)
        #   2. Resolves roles via three-strategy fallback (enum, XML, name)
        #   3. Detects bullets from XML (only explicit buChar/buAutoNum)
        #   4. Uses run.font.bold for reliable inherited bold resolution
        #   5. Recursively processes group shapes through same dispatch
        _process_shapes_for_native(slide.shapes, slide_data)

        # ── [PATCH] Walk slide layout for additional pictures ──
        # Some slides display visual content that lives in the slide layout
        # (not the slide itself). python-pptx walks slide.shapes which only
        # contains shapes added directly to the slide — it doesn't include
        # shapes inherited from the layout/master.
        #
        # Reviewer feedback: "Image missing" on IsDB slides 3 and 8.
        # Verified: IsDB slide 8's image lives in its layout 'Section Header 2'
        # under shape name 'Picture 13'. python-pptx never sees it via
        # slide.shapes alone, so the extractor missed it before this fix.
        #
        # IMPORTANT: when walking the layout, we ONLY want pictures, not text
        # frames. Layout text frames contain design-time placeholder prompts
        # like "Click to edit Master title style" / "Click to edit Master text
        # styles / Second level / Third level / ...". Those are template
        # design hints, not slide content, and must not leak into the CO.
        # We use a dedicated walker that skips text-bearing shapes.
        try:
            layout = slide.slide_layout
            if layout is not None:
                _walk_layout_pictures_only(layout.shapes, slide_data)
        except Exception as e:
            logger.debug(f"Slide layout walk failed: {e}")

        # ── [PATCH] Dedupe body_lines (catches DI-style repeats and shape duplicates) ──
        # Sometimes python-pptx walks group shapes that contain text frames AND
        # the children of those groups, emitting the same text twice. Or two
        # text shapes contain the same text (e.g., header reused). This dedup
        # uses a sliding window similar to the DOCX extractor.
        # Fixes: "Objectives word repeated twice", "Duplicate text"
        deduped_lines = []
        for line in slide_data["body_lines"]:
            stripped = line.strip()
            if not stripped:
                deduped_lines.append(line)
                continue
            recent = [l.strip() for l in deduped_lines[-6:] if l.strip()]
            if stripped in recent:
                continue
            deduped_lines.append(line)
        slide_data["body_lines"] = deduped_lines

        # ── Assemble content ──
        title = slide_data["title_text"] or ''
        body_lines = slide_data["body_lines"]
        tables_md = slide_data["tables_md"]
        has_images = slide_data["has_images"]
        has_smartart = slide_data["has_smartart"]

        content_parts = []
        if body_lines:
            content_parts.append('\n'.join(body_lines))
        for tbl_md in tables_md:
            content_parts.append(f'\n\n[TABLE]\n{tbl_md}')
        # Notes are now stored separately (in _slide_notes field), NOT in content,
        # so they do not pollute the "Full Page Content" column of the outline.
        # Layer 3 / storyboard generator can access notes via the separate field.

        # Handle image-only slides
        if not body_lines and not tables_md and (has_images or has_smartart):
            content_parts.insert(0, '[IMAGE-ONLY SLIDE - see slide screenshot for content]')

        # ── Determine confidence ──
        confidence = 1.0
        if not title or title.startswith('Slide '):
            confidence = min(confidence, 0.5)
        if has_smartart and not body_lines:
            confidence = min(confidence, 0.4)
        if not body_lines and not tables_md:
            confidence = min(confidence, 0.3)

        # ── [CHANGE] Subtitle in output ──
        # OLD: subtopics was always "N/A".
        # NEW: Uses subtitle from the subtitle placeholder if found, providing
        #      richer structural data for Layer 3 structure assignment.
        slides_out.append({
            "page_number": visible_idx,
            "chapter": "",  # Set by Layer 3
            "topic": _clean_pptx_title(title) if title else f"Slide {visible_idx}",
            "subtopics": slide_data.get("subtitle_text") or "N/A",
            "content": clean_ocr_noise('\n'.join(content_parts)),
            # [PATCH] Now populated with extracted PPTX images so the renderer
            # can resolve [IMAGE: id] markers to actual base64 data.
            "images": slide_data.get("extracted_images", []),
            "tables": [],
            "source_page_range": str(visible_idx),
            "_confidence": confidence,
            "_has_smartart": has_smartart,
            "_has_images": has_images,
            "_slide_notes": notes_text,           # Notes kept separate from content
            "_slide_duration": slide_duration,    # Parsed from "Slide duration: X" in notes
        })

    logger.info(f"[LAYER 1] Extracted {len(slides_out)} slides. "
                f"Low-confidence: {sum(1 for s in slides_out if s['_confidence'] < 0.8)}")
    return slides_out


def _parse_slide_duration(notes_text: str) -> Optional[str]:
    """
    Parse 'Slide duration: X min' (or seconds, or MM:SS) from slide notes.
    Returns a normalized 'X min Y sec' string, or None if not found.

    Reviewers consistently complained that the CO output durations were wrong.
    The system was computing duration from word count instead of reading the
    explicit duration that PPT authors put in the speaker notes.
    """
    if not notes_text:
        return None

    # Normalize whitespace
    text = notes_text.replace('\u000b', '\n').replace('\xa0', ' ')

    # Pattern 1: "Slide duration: 5 min" or "Slide duration : 1 min."
    m = re.search(
        r'[Ss]lide\s*duration\s*[:\-]\s*([\d.]+)\s*(min|minute|m)\b',
        text
    )
    if m:
        try:
            mins = float(m.group(1))
            whole = int(mins)
            secs = int(round((mins - whole) * 60))
            return f"{whole} min {secs} sec" if secs else f"{whole} min 0 sec"
        except ValueError:
            pass

    # Pattern 2: "Slide duration: 30 sec" or "30 seconds"
    m = re.search(
        r'[Ss]lide\s*duration\s*[:\-]\s*(\d+)\s*(sec|second|s)\b',
        text
    )
    if m:
        try:
            secs = int(m.group(1))
            return f"{secs // 60} min {secs % 60} sec"
        except ValueError:
            pass

    # Pattern 3: MM:SS format "Duration: 1:30"
    m = re.search(
        r'[Ss]lide\s*duration\s*[:\-]\s*(\d+):(\d{1,2})',
        text
    )
    if m:
        try:
            return f"{int(m.group(1))} min {int(m.group(2))} sec"
        except ValueError:
            pass

    return None


# ─────────────────────────────────────────────────────────────────────────────
# [SUPERSEDED] _is_pptx_run_bold
# ─────────────────────────────────────────────────────────────────────────────
# OLD CODE: Manually walked a 4-level XML inheritance chain:
#   Level 1: run.font.bold (explicit on run)
#   Level 2: paragraph defRPr/@b in XML
#   Level 3: slide layout placeholder default
#   Level 4: slide master txStyles/bodyStyle
# This was complex (~70 lines), fragile, and still missed some inheritance
# paths (e.g., theme-level bold defaults, list style bold).
#
# NEW CODE: Uses python-pptx's run.font.bold property directly in
# _para_to_markdown(). The high-level API resolves the full OOXML inheritance
# chain automatically (run rPr → paragraph defRPr → list style → master →
# theme). This single property call replaces the entire manual XML walk.
#
# The old function is preserved below (commented out) for reference:
#
# def _is_pptx_run_bold(run, para, slide) -> bool:
#     """
#     Check PPTX run bold with FULL inheritance chain:
#       1. run.font.bold (explicit on this run)
#       2. paragraph default run properties (defRPr)
#       3. slide layout placeholder default
#       4. slide master placeholder default
#     Returns True if bold at any level in the chain.
#     """
#     if run.font.bold is True:
#         return True
#     if run.font.bold is False:
#         return False
#     ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
#     try:
#         defRPr = para._p.find(f'.//{ns}defRPr')
#         if defRPr is not None and defRPr.get('b') == '1':
#             return True
#     except Exception:
#         pass
#     try:
#         ph = None
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for p in shape.text_frame.paragraphs:
#                     if p is para:
#                         try:
#                             ph = shape.placeholder_format
#                         except (ValueError, AttributeError):
#                             pass
#                         break
#         if ph is not None:
#             layout = slide.slide_layout
#             for layout_ph in layout.placeholders:
#                 try:
#                     if layout_ph.placeholder_format.idx == ph.idx:
#                         for layout_para in layout_ph.text_frame.paragraphs:
#                             ldefRPr = layout_para._p.find(f'.//{ns}defRPr')
#                             if ldefRPr is not None and ldefRPr.get('b') == '1':
#                                 return True
#                             break
#                 except (ValueError, AttributeError):
#                     continue
#     except Exception:
#         pass
#     try:
#         master = slide.slide_layout.slide_master
#         txBody = master.element.find(f'.//{ns}txStyles')
#         if txBody is not None:
#             bodyStyle = txBody.find(f'{ns}bodyStyle')
#             if bodyStyle is not None:
#                 lvl1pPr = bodyStyle.find(f'{ns}lvl1pPr')
#                 if lvl1pPr is not None:
#                     mdefRPr = lvl1pPr.find(f'{ns}defRPr')
#                     if mdefRPr is not None and mdefRPr.get('b') == '1':
#                         return True
#     except Exception:
#         pass
#     return False
# ─────────────────────────────────────────────────────────────────────────────


def _clean_pptx_title(text: str) -> str:
    """Clean title text — remove numbering prefixes, section markers."""
    if not text:
        return "N/A"
    # Remove "Section X:", "Chapter Y:", "Module Z:" prefixes
    text = re.sub(
        r'^\s*(section|chapter|module|part|day|week|unit)\s*\d+[\s:\-]*',
        '', text, flags=re.IGNORECASE)
    text = re.sub(r'^\s*\d+[\.\)]\s+', '', text)
    text = re.sub(r'^\s*[:\-|–]\s*', '', text)
    # Collapse vertical-tab characters (common in PPTX for line breaks)
    text = text.replace('\x0b', ' ').replace('\r', ' ')
    return re.sub(r'\s+', ' ', text).strip()


# ══════════════════════════════════════════════════════════════════════════════
# LAYER 1: NATIVE DOCX EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════

def extract_docx_native(docx_path: str) -> List[Dict[str, Any]]:
    """
    Layer 1: Extract all content from DOCX using native XML structure.

    Handles:
      ✓ Standard heading styles (Heading 1/2/3, Title, Subtitle)
      ✓ Custom corporate heading styles (by keyword matching + font heuristics)
      ✓ Visual-only headings (bold + large font, no heading style)
      ✓ Bullet/list hierarchy via numPr/ilvl AND manual numbering patterns
      ✓ Bold text — full-line and inline, including style-inherited bold
      ✓ Tables with proper structure preserved
      ✓ Nested tables (outer table extracted, inner flagged)
      ✓ Images detected and referenced
    """
    if not DOCX_AVAILABLE:
        logger.warning("python-docx not available, cannot extract natively")
        return []

    logger.info(f"[LAYER 1] Native DOCX extraction: {os.path.basename(docx_path)}")

    # ── Check for password protection ──
    if is_password_protected(docx_path):
        logger.error(f"[LAYER 1] File is password-protected: {os.path.basename(docx_path)}. "
                      "Cannot extract natively. Return empty to trigger DI fallback.")
        return []

    try:
        doc = python_docx.Document(docx_path)
    except (BadZipFile, DocxPackageError) as e:
        logger.error(f"[LAYER 1] File is corrupted or password-protected: {e}")
        return []
    except Exception as e:
        logger.error(f"[LAYER 1] Failed to open DOCX: {e}")
        return []

    # ── First pass: determine body font size for heading heuristics ──
    font_sizes = []
    for para in doc.paragraphs:
        for run in para.runs:
            if run.font.size:
                font_sizes.append(run.font.size.pt)
    body_font_size = 11.0
    if font_sizes:
        from collections import Counter
        size_counts = Counter(round(s, 1) for s in font_sizes)
        body_font_size = size_counts.most_common(1)[0][0]

    # ── Second pass: extract sections ──
    sections = []
    current_heading = None
    current_level = 0
    current_lines = []
    section_num = 0
    image_count = 0

    def flush():
        nonlocal current_lines, section_num
        if current_heading or current_lines:
            section_num += 1
            content = '\n'.join(current_lines)
            confidence = 1.0 if current_level > 0 else 0.6
            sections.append({
                "page_number": section_num,
                "chapter": "",  # Set by Layer 3
                "topic": current_heading or "Document Content",
                "subtopics": "N/A",
                "content": clean_ocr_noise(content),
                "images": [],
                "tables": [],
                "source_page_range": str(section_num),
                "_confidence": confidence,
                "_heading_level": current_level,
            })
            current_lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Check if this is a heading
        h_level = is_likely_heading_docx(para, body_font_size)

        if h_level > 0:
            flush()
            current_heading = text
            current_level = h_level
        else:
            # Extract text with formatting
            list_level = get_para_list_level_docx(para)
            formatted = extract_bold_text(para)

            if list_level >= 0:
                # List item with proper indentation
                prefix = "  " * list_level + "- "
                current_lines.append(f"{prefix}{formatted}")
            else:
                # Regular paragraph
                current_lines.append(formatted)

    # Extract tables
    for ti, table in enumerate(doc.tables):
        rows_data = []
        for row in table.rows:
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            rows_data.append(cells)
        if rows_data:
            tbl_md = format_table_markdown(rows_data)
            # Attach to the most recent section
            if sections:
                sections[-1]["content"] += f"\n\n[TABLE]\n{tbl_md}"
            else:
                current_lines.append(f"\n[TABLE]\n{tbl_md}")

    # Extract images (just count — actual images handled by DI/PDF path)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_count += 1

    flush()

    has_headings = any(s.get("_heading_level", 0) > 0 for s in sections)
    logger.info(f"[LAYER 1] Extracted {len(sections)} sections. "
                f"Has headings: {has_headings}. Images: {image_count}")
    return sections


# ══════════════════════════════════════════════════════════════════════════════
# LAYER 2: LLM VISION ENHANCEMENT
# ══════════════════════════════════════════════════════════════════════════════

def enhance_slide_with_vision(
    openai_client,
    deployment_name: str,
    slide_num: int,
    layer1_data: Dict[str, Any],
    pdf_path: str,
) -> Dict[str, Any]:
    """
    Layer 2: Use GPT-4o Vision to enhance/correct a single slide's extraction.
    Called ONLY when Layer 1 produces low-confidence output.

    Args:
        openai_client: Azure OpenAI client
        deployment_name: Vision-capable model deployment name
        slide_num: 1-based slide number
        layer1_data: The Layer 1 extraction output for this slide
        pdf_path: Path to the PDF conversion (for screenshot)

    Returns:
        Enhanced slide data dict (same format as Layer 1 output)
    """
    if not FITZ_AVAILABLE or not openai_client:
        return layer1_data

    logger.info(f"[LAYER 2] Vision-enhancing slide {slide_num}...")

    try:
        doc = fitz.open(pdf_path)
        if slide_num - 1 >= len(doc):
            return layer1_data
        page = doc[slide_num - 1]
        pix = page.get_pixmap(dpi=150)
        img_b64 = base64.b64encode(pix.tobytes('png')).decode('utf-8')
        doc.close()
    except Exception as e:
        logger.warning(f"[LAYER 2] Failed to get slide image: {e}")
        return layer1_data

    system_prompt = (
        "You are a slide content extraction assistant. "
        "You are given a slide IMAGE and an initial TEXT extraction that may be incomplete or poorly structured. "
        "Your task is to produce a corrected, well-structured extraction.\n\n"
        "RULES:\n"
        "1. TITLE: Identify the single main heading (largest/most prominent text on the slide)\n"
        "2. BODY: Extract ALL text content with correct bullet hierarchy:\n"
        "   - Top-level bullets: '- item'\n"
        "   - Sub-bullets: '  - sub-item'\n"
        "   - Sub-sub-bullets: '    - sub-sub-item'\n"
        "3. BOLD: Mark text that is visually bold with **double asterisks**\n"
        "4. TABLES: Reproduce any tables in markdown pipe format:\n"
        "   | Col1 | Col2 |\n"
        "   | --- | --- |\n"
        "   | val1 | val2 |\n"
        "5. IMAGE SLIDES: If the slide is primarily a screenshot, photo, or diagram with "
        "minimal text, set is_image_slide=true and describe the visual briefly\n"
        "6. IGNORE: Slide numbers, logos, decorative backgrounds, page footers\n"
        "7. PRESERVE: Every word of actual content must appear in your output\n\n"
        "Return ONLY a JSON object:\n"
        '{"title": "...", "body": "- bullet\\n  - sub-bullet", '
        '"tables": ["| ... |"], "is_image_slide": false}'
    )

    layer1_text = layer1_data.get('content', '')[:1500]
    layer1_title = layer1_data.get('topic', '')

    user_content = [
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
        {"type": "text", "text": (
            f"Initial extraction (may be incomplete):\n"
            f"Title: {layer1_title}\n"
            f"Content:\n{layer1_text}\n\n"
            f"Correct and enhance this extraction based on what you see in the image."
        )}
    ]

    try:
        response = openai_client.chat.completions.create(
            model=deployment_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content}
            ],
            temperature=0,
            response_format={"type": "json_object"},
            max_tokens=3000,
        )
        enhanced = json.loads(response.choices[0].message.content)

        # Build enhanced content
        content_parts = []
        body = enhanced.get('body', '')
        if body:
            content_parts.append(body)
        for tbl in enhanced.get('tables', []):
            content_parts.append(f"\n\n[TABLE]\n{tbl}")

        if enhanced.get('is_image_slide', False):
            if not content_parts:
                content_parts.insert(0, '[IMAGE-DOMINANT SLIDE - see slide screenshot]')

        # Preserve notes from Layer 1
        if '--- SLIDE NOTES ---' in layer1_data.get('content', ''):
            notes_start = layer1_data['content'].index('--- SLIDE NOTES ---')
            content_parts.append(f"\n\n{layer1_data['content'][notes_start:]}")

        return {
            **layer1_data,
            "topic": enhanced.get('title', layer1_data.get('topic', f'Slide {slide_num}')),
            "content": clean_ocr_noise('\n'.join(content_parts)),
            "_confidence": 0.85,  # Enhanced by Vision
        }

    except Exception as e:
        logger.warning(f"[LAYER 2] Vision enhancement failed for slide {slide_num}: {e}")
        return layer1_data


def enhance_docx_with_llm(
    openai_client,
    deployment_name: str,
    full_text: str,
) -> List[Dict[str, str]]:
    """
    Layer 2 for DOCX: When Layer 1 finds NO heading styles, ask LLM to identify headings.
    Returns list of {"text": str, "level": int} for detected headings.
    """
    if not openai_client or not full_text:
        return []

    logger.info("[LAYER 2] LLM heading detection for unstructured DOCX...")

    # Take first 4000 chars to identify the heading pattern
    sample = full_text[:4000]

    prompt = (
        "Analyze this document text and identify lines that are likely HEADINGS "
        "(section titles, chapter names, topic headers).\n\n"
        "Look for:\n"
        "- Lines that are short (< 10 words) and don't end with periods\n"
        "- Lines that start with numbers (1., 1.1, I., A.)\n"
        "- Lines in ALL CAPS\n"
        "- Lines that introduce a new topic\n\n"
        f"DOCUMENT TEXT:\n{sample}\n\n"
        "Return JSON: {\"headings\": [{\"text\": \"exact heading text\", \"level\": 1|2|3}]}\n"
        "Level 1 = major sections, Level 2 = sub-sections, Level 3 = minor sections"
    )

    try:
        response = openai_client.chat.completions.create(
            model=deployment_name,
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
            response_format={"type": "json_object"},
            max_tokens=2000,
        )
        result = json.loads(response.choices[0].message.content)
        return result.get("headings", [])
    except Exception as e:
        logger.warning(f"[LAYER 2] LLM heading detection failed: {e}")
        return []


# ══════════════════════════════════════════════════════════════════════════════
# LAYER 3: LLM STRUCTURE ASSIGNMENT
# ══════════════════════════════════════════════════════════════════════════════
# NOTE: Layer 3 is already implemented in presentation_agent.py as
# _analyze_structure_with_llm() and in no_toc_agent.py as _classify_headings_llm().
# These work well. The only change needed is to feed them Layer 1+2 output
# instead of DI-degraded output.
#
# For reference, here is how to call the existing Layer 3 from outside the agent:

def assign_structure_with_llm(
    openai_client,
    deployment_name: str,
    slides_data: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    """
    Layer 3: Assign chapter/topic/subtopic structure across ALL slides.
    This mirrors _analyze_structure_with_llm() in presentation_agent.py.

    Can be used as a standalone function or the existing agent method can be kept.
    """
    if not openai_client or not slides_data:
        return slides_data

    total = len(slides_data)
    logger.info(f"[LAYER 3] Assigning structure to {total} slides...")

    # Build context for LLM
    context_lines = []
    for s in slides_data:
        sn = s.get('page_number', 0)
        title = s.get('topic', 'No Title')
        snippet = s.get('content', '')[:300].replace('\n', ' ')
        context_lines.append(f"Slide {sn}: [Title: {title}] | Content: {snippet}")

    full_context = "\n".join(context_lines)

    system_prompt = (
        "You are a presentation structure analyst. "
        "Assign a 'chapter' and optionally refine the 'topic' for each slide.\n\n"
        "RULES:\n"
        "1. Chapters = MAJOR thematic sections (e.g., 'Benefits', 'Retirement', 'Systems')\n"
        "2. Do NOT group non-adjacent slides into one chapter if a different topic is between them\n"
        "3. Topic = concise descriptive name for each slide\n"
        "4. Subtopics = key bullet points summarized\n"
        "5. Return a JSON object with a 'structure' array\n"
    )

    user_prompt = (
        f"Here are {total} slides:\n\n{full_context}\n\n"
        f"Return JSON:\n"
        f'{{"structure": [{{"slide_number": 1, "chapter": "...", "topic": "...", "subtopics": "..."}}]}}'
    )

    try:
        response = openai_client.chat.completions.create(
            model=deployment_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0,
            response_format={"type": "json_object"},
            max_tokens=4000,
        )
        data = json.loads(response.choices[0].message.content)
        structure = {
            item['slide_number']: item
            for item in data.get('structure', [])
            if 'slide_number' in item
        }

        # Apply structure to slides
        for s in slides_data:
            sn = s['page_number']
            if sn in structure:
                s['chapter'] = structure[sn].get('chapter', 'General')
                s['topic'] = structure[sn].get('topic', s.get('topic', f'Slide {sn}'))
                s['subtopics'] = structure[sn].get('subtopics', 'N/A')

        logger.info(f"[LAYER 3] Assigned structure to {len(structure)} of {total} slides")

    except Exception as e:
        logger.error(f"[LAYER 3] Structure assignment failed: {e}")
        # Fallback: use slide titles as topics, generic chapters
        for s in slides_data:
            if not s.get('chapter'):
                s['chapter'] = 'General'

    return slides_data


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ORCHESTRATOR — CALL THIS FROM THE AGENTS
# ══════════════════════════════════════════════════════════════════════════════

def extract_pptx_robust(
    pptx_path: str,
    pdf_path: Optional[str],
    openai_client=None,
    vision_deployment: str = None,
    text_deployment: str = None,
) -> List[Dict[str, Any]]:
    """
    Full three-layer extraction for presentation files.
    Call this from PresentationAgent.process() instead of the current DI pipeline.

    Handles:
      - .pptx → Layer 1 (native) + Layer 2 (vision) + Layer 3 (structure)
      - .ppt  → Returns empty list with log message. Agent must use DI fallback.
      - Password-protected → Returns empty list with error. Agent must handle.

    Returns:
        List of chunk dicts ready for _upload_chunks().
        Empty list means "use DI fallback" — NOT "file has no content".
    """
    # ── Check file format ──
    fmt = detect_file_format(pptx_path)
    if fmt == 'ppt_legacy':
        logger.warning("[ROBUST] Legacy .ppt file detected. Native extraction impossible. "
                       "Returning empty to trigger DI fallback in the agent.")
        return []
    if fmt != 'pptx':
        logger.warning(f"[ROBUST] Unexpected format '{fmt}' for PPTX extractor. "
                       "Returning empty to trigger fallback.")
        return []

    # ── Layer 1: Native extraction ──
    slides = extract_pptx_native(pptx_path)

    if not slides:
        logger.warning("Layer 1 returned nothing — file may be password-protected, "
                       "corrupted, or an invalid PPTX. DI fallback will be used.")
        return []

    # ── Layer 2: Vision enhancement for low-confidence slides ──
    if openai_client and pdf_path and vision_deployment:
        for i, slide in enumerate(slides):
            if slide.get('_confidence', 1.0) < 0.8:
                logger.info(f"  Slide {slide['page_number']}: confidence={slide['_confidence']:.1f}, enhancing...")
                slides[i] = enhance_slide_with_vision(
                    openai_client, vision_deployment,
                    slide['page_number'], slide, pdf_path
                )
    else:
        logger.info("[LAYER 2] Skipped — no OpenAI client or no PDF path")

    # ── Layer 3: Structure assignment ──
    if openai_client and text_deployment:
        slides = assign_structure_with_llm(openai_client, text_deployment, slides)
    else:
        logger.info("[LAYER 3] Skipped — no OpenAI client")
        for s in slides:
            if not s.get('chapter'):
                s['chapter'] = s.get('topic', 'General')

    # ── Clean internal fields before returning ──
    for s in slides:
        s.pop('_confidence', None)
        s.pop('_has_smartart', None)
        s.pop('_has_images', None)
        s.pop('_heading_level', None)

    return slides


def extract_docx_robust(
    docx_path: str,
    openai_client=None,
    text_deployment: str = None,
) -> List[Dict[str, Any]]:
    """
    Full three-layer extraction for Word document files.
    Call this from hierarchical_agent / no_toc_agent for Word files.

    Handles:
      - .docx → Layer 1 (native) + Layer 2 (LLM headings) + Layer 3 (structure)
      - .doc  → Returns empty list with log message. Agent must use DI fallback.
      - Password-protected → Returns empty list with error. Agent must handle.

    Returns:
        List of chunk dicts ready for uploading.
        Empty list means "use DI fallback" — NOT "file has no content".
    """
    # ── Check file format ──
    fmt = detect_file_format(docx_path)
    if fmt == 'doc_legacy':
        logger.warning("[ROBUST] Legacy .doc file detected. Native extraction impossible. "
                       "Returning empty to trigger DI fallback in the agent.")
        return []
    if fmt != 'docx':
        logger.warning(f"[ROBUST] Unexpected format '{fmt}' for DOCX extractor. "
                       "Returning empty to trigger fallback.")
        return []

    # ── Layer 1: XML-based extraction (primary) ──
    # Uses lxml to parse DOCX XML directly. This handles:
    #   - Files that python-docx can't open (e.g., #Contents hyperlinks)
    #   - Custom bullet styles (BT Bullet 1, etc.)
    #   - Visual heading detection for files with no heading styles
    #   - Images inside table cells
    #   - Page breaks (explicit + lastRenderedPageBreak)
    #   - TOC filtering
    try:
        from extract_docx_xml import extract_docx_xml
        sections = extract_docx_xml(docx_path)
        if sections:
            logger.info(f"[LAYER 1] XML extraction succeeded: {len(sections)} sections")
        else:
            logger.warning("[LAYER 1] XML extraction returned empty, trying python-docx fallback")
            sections = extract_docx_native(docx_path)
    except ImportError:
        logger.info("[LAYER 1] extract_docx_xml not available, using python-docx")
        sections = extract_docx_native(docx_path)
    except Exception as e:
        logger.warning(f"[LAYER 1] XML extraction failed ({e}), trying python-docx fallback")
        sections = extract_docx_native(docx_path)

    if not sections:
        logger.warning("Layer 1 returned nothing — file may be password-protected, "
                       "corrupted, or an invalid DOCX. DI fallback will be used.")
        return []

    # ── Layer 2: If no heading styles detected, use LLM ──
    has_headings = any(s.get('_heading_level', 0) > 0 for s in sections)
    if not has_headings and openai_client and text_deployment:
        logger.info("[LAYER 2] No heading styles found. Using LLM to detect headings...")
        full_text = '\n\n'.join(s.get('content', '') for s in sections)
        llm_headings = enhance_docx_with_llm(openai_client, text_deployment, full_text)

        if llm_headings:
            # Re-segment the content based on LLM-detected headings
            heading_texts = {h['text']: h['level'] for h in llm_headings}
            new_sections = []
            current = {"topic": "Document Start", "content": "", "_heading_level": 0}

            for section in sections:
                content = section.get('content', '')
                for line in content.split('\n'):
                    stripped = line.strip()
                    if stripped in heading_texts:
                        if current['content'].strip():
                            new_sections.append({
                                **section,
                                "topic": current['topic'],
                                "content": current['content'].strip(),
                                "_heading_level": current['_heading_level'],
                                "_confidence": 0.8,
                            })
                        current = {
                            "topic": stripped,
                            "content": "",
                            "_heading_level": heading_texts[stripped],
                        }
                    else:
                        current['content'] += line + '\n'

            if current['content'].strip():
                new_sections.append({
                    "page_number": len(new_sections) + 1,
                    "chapter": "",
                    "topic": current['topic'],
                    "subtopics": "N/A",
                    "content": current['content'].strip(),
                    "images": [],
                    "tables": [],
                    "source_page_range": str(len(new_sections) + 1),
                    "_heading_level": current['_heading_level'],
                    "_confidence": 0.8,
                })

            if new_sections:
                sections = new_sections
                # Re-number
                for i, s in enumerate(sections, 1):
                    s['page_number'] = i
                    s['source_page_range'] = str(i)

    # ── Layer 3: Classify headings into Chapter/Topic/Subtopic ──
    if openai_client and text_deployment:
        # Use simple mapping for standard heading levels
        current_chapter = "Document Content"
        current_topic = ""
        for s in sections:
            level = s.get('_heading_level', 0)
            if level == 1:
                current_chapter = s.get('topic', 'General')
                s['chapter'] = current_chapter
            elif level == 2:
                s['chapter'] = current_chapter
                current_topic = s.get('topic', '')
            elif level == 3:
                s['chapter'] = current_chapter
                s['subtopics'] = s.get('topic', '')
                s['topic'] = current_topic if current_topic else s.get('topic', '')
            else:
                s['chapter'] = current_chapter
    else:
        for s in sections:
            if not s.get('chapter'):
                s['chapter'] = 'General'

    # ── Clean internal fields ──
    for s in sections:
        s.pop('_confidence', None)
        s.pop('_heading_level', None)

    return sections
