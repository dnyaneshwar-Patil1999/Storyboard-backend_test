import argparse
import json
import base64
import os
import sys
import subprocess
import shutil
import re
import time
import traceback
from collections import defaultdict
from abc import ABC, abstractmethod
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError

# --- IMPORTS ---
from azure.core.credentials import AzureKeyCredential
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.ai.documentintelligence.models import AnalyzeDocumentRequest, DocumentContentFormat
from azure.storage.blob import BlobServiceClient
from openai import AzureOpenAI
import fitz  # PyMuPDF
from pptx import Presentation

# --- ENVIRONMENT CHECKER ---

def check_function_app_environment():
    """Check if running in Azure Function App and log environment details"""
    print("=" * 80)
    print("ENVIRONMENT DIAGNOSTICS")
    print("=" * 80)
    
    print(f"[DEBUG] Python version: {sys.version}")
    print(f"[DEBUG] Platform: {sys.platform}")
    
    # Check Azure Function App environment variables
    azure_env_vars = [
        'WEBSITE_SITE_NAME',
        'FUNCTIONS_WORKER_RUNTIME',
        'FUNCTIONS_EXTENSION_VERSION',
        'AzureWebJobsStorage'
    ]
    
    for var in azure_env_vars:
        value = os.environ.get(var, 'Not set')
        print(f"[DEBUG] {var}: {value}")
    
    # Check required environment variables
    required_vars = [
        'AZURE_OPENAI_ENDPOINT',
        'AZURE_OPENAI_API_KEY',
        'AZURE_DI_ENDPOINT',
        'AZURE_DI_KEY'
    ]
    
    print("\n[DEBUG] Required Environment Variables:")
    for var in required_vars:
        exists = var in os.environ
        status = "✓" if exists else "✗"
        if exists:
            # Don't print full keys for security
            if 'KEY' in var or 'SECRET' in var:
                value = os.environ[var][:10] + "..." if len(os.environ[var]) > 10 else "***"
            else:
                value = os.environ[var]
            print(f"  {status} {var}: {value}")
        else:
            print(f"  {status} {var}: Missing")
    
    print("=" * 80)

# --- BASE AGENT CLASS ---

class BaseAgent(ABC):
    def __init__(self, config, local_doc_path, blob_path, toc_data, toc_page_range=None):
        self.local_doc_path = local_doc_path
        self.blob_path = blob_path
        self.config = config
        self.toc_data = toc_data
        self.toc_page_range = set(p + 1 for p in (toc_page_range or []))
        
        print(f"[DEBUG] BaseAgent initialized with:")
        print(f"[DEBUG]   Local path: {local_doc_path}")
        print(f"[DEBUG]   Blob path: {blob_path}")
        print(f"[DEBUG]   Config keys: {list(config.keys())}")

        # Initialize Document Intelligence Client
        try:
            print("[DEBUG] Initializing Document Intelligence Client...")
            self.document_intelligence_client = DocumentIntelligenceClient(
                endpoint=config["AZURE_DI_ENDPOINT"],
                credential=AzureKeyCredential(config["AZURE_DI_KEY"])
            )
            print("[DEBUG] Document Intelligence Client initialized successfully")
        except Exception as e:
            print(f"[ERROR] Failed to initialize Document Intelligence Client: {e}")
            raise

        # Initialize Azure OpenAI Client - USING SAME PATTERN AS HIERARCHICALCHUNKINGAGENT
        print("[DEBUG] Checking OpenAI configuration...")
        self.openai_client = None
        self.openai_deployment = config.get("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o")
        
        # Check if OpenAI config exists in config dict (for backward compatibility)
        if "AZURE_OPENAI_ENDPOINT" in config and "AZURE_OPENAI_API_KEY" in config:
            if config["AZURE_OPENAI_ENDPOINT"] and config["AZURE_OPENAI_API_KEY"]:
                try:
                    print("[DEBUG] Initializing Azure OpenAI Client...")
                    self.openai_client = AzureOpenAI(
                        azure_endpoint=config["AZURE_OPENAI_ENDPOINT"],
                        api_key=config["AZURE_OPENAI_API_KEY"],
                        api_version="2024-02-01"
                    )
                    print(f"[DEBUG] Azure OpenAI Client initialized successfully with deployment: {self.openai_deployment}")
                except Exception as e:
                    print(f"[ERROR] Failed to initialize Azure OpenAI Client: {e}")
                    print("[WARN] LLM functionality will be disabled")
            else:
                print("[WARN] OpenAI endpoint or key is empty in config")
        # Also check environment variables (for Function App)
        elif os.environ.get("AZURE_OPENAI_ENDPOINT") and os.environ.get("AZURE_OPENAI_API_KEY"):
            print("[DEBUG] Using OpenAI config from environment variables...")
            try:
                self.openai_client = AzureOpenAI(
                    azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
                    api_key=os.environ["AZURE_OPENAI_API_KEY"],
                    api_version="2024-02-01"
                )
                print(f"[DEBUG] Azure OpenAI Client initialized from environment variables")
            except Exception as e:
                print(f"[ERROR] Failed to initialize Azure OpenAI Client from environment: {e}")
        else:
            print("[WARN] OpenAI configuration not found in config or environment")
            print(f"[DEBUG] Config has AZURE_OPENAI_ENDPOINT: {'AZURE_OPENAI_ENDPOINT' in config}")
            print(f"[DEBUG] Config has AZURE_OPENAI_API_KEY: {'AZURE_OPENAI_API_KEY' in config}")
            print(f"[DEBUG] Env has AZURE_OPENAI_ENDPOINT: {'AZURE_OPENAI_ENDPOINT' in os.environ}")
            print(f"[DEBUG] Env has AZURE_OPENAI_API_KEY: {'AZURE_OPENAI_API_KEY' in os.environ}")
        
        # Initialize Blob Storage Client
        try:
            print("[DEBUG] Initializing Blob Storage Client...")
            self.dest_container_client = BlobServiceClient.from_connection_string(
                config["STORAGE_CONNECTION_STRING"]
            ).get_container_client(config["DESTINATION_CONTAINER_NAME"])
            print("[DEBUG] Blob Storage Client initialized successfully")
        except Exception as e:
            print(f"[ERROR] Failed to initialize Blob Storage Client: {e}")
            raise

    @abstractmethod
    def process(self): 
        pass

    def _convert_to_pdf(self):
        """Common PDF conversion logic using LibreOffice."""
        if self.local_doc_path.lower().endswith('.pdf'):
            print(f"[INFO] File is already PDF: {os.path.basename(self.local_doc_path)}")
            return self.local_doc_path

        print("[INFO] Converting file to PDF...")
        try:
            pdf_path = os.path.splitext(self.local_doc_path)[0] + "_processing.pdf"
            print(f"[DEBUG] Target PDF path: {pdf_path}")
            
            # Check if LibreOffice is available
            try:
                subprocess.run(["libreoffice", "--version"], capture_output=True, check=True)
                print("[DEBUG] LibreOffice is available")
            except FileNotFoundError:
                raise RuntimeError("`libreoffice` command not found. Please ensure LibreOffice is installed.")
            
            # Run LibreOffice headless conversion
            print(f"[DEBUG] Running LibreOffice conversion...")
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir",
                 os.path.dirname(self.local_doc_path), self.local_doc_path],
                check=True, timeout=300, capture_output=True, text=True
            )
            print(f"[DEBUG] LibreOffice stdout: {result.stdout[:200]}")
            if result.stderr:
                print(f"[DEBUG] LibreOffice stderr: {result.stderr[:200]}")
            
            # LibreOffice output naming handling
            base_name = os.path.splitext(os.path.basename(self.local_doc_path))[0]
            expected_pdf = os.path.join(os.path.dirname(self.local_doc_path), f"{base_name}.pdf")
            print(f"[DEBUG] Expected PDF: {expected_pdf}")
            print(f"[DEBUG] File exists: {os.path.exists(expected_pdf)}")
            
            if not os.path.exists(expected_pdf):
                # List files in directory for debugging
                dir_files = os.listdir(os.path.dirname(self.local_doc_path))
                print(f"[DEBUG] Files in directory: {dir_files}")
                raise RuntimeError(f"PDF conversion failed to create the expected file: {expected_pdf}")
            
            if os.path.abspath(expected_pdf) != os.path.abspath(pdf_path):
                print(f"[DEBUG] Moving {expected_pdf} to {pdf_path}")
                shutil.move(expected_pdf, pdf_path)
            else:
                print(f"[DEBUG] PDF already at correct location")
                
            return pdf_path
            
        except subprocess.TimeoutExpired:
            raise RuntimeError("PDF conversion timed out after 300 seconds")
        except subprocess.CalledProcessError as e:
            raise RuntimeError(f"PDF conversion failed with return code {e.returncode}: {e.stderr}")
        except Exception as e:
            raise RuntimeError(f"PDF conversion failed: {str(e)}")

    def _upload_chunks(self, data):
        print(f"\n[INFO] Uploading {len(data)} chunks to Azure...")
        bp, fn = os.path.dirname(self.blob_path), os.path.basename(self.blob_path)
        md = self._extract_metadata()

        for i, chunk_data in enumerate(data):
            chunk_num = i + 1
            
            # Determine logic for filename parts
            path_parts = [p for p in [chunk_data.get('chapter'), chunk_data.get('topic')] if p and p != 'N/A']
            
            filename_prefix_parts = [p for p in [md.get('Client'), md.get('Project'), md.get('Module')] if p and p != 'N/A']
            filename_prefix = "-".join(filename_prefix_parts) if filename_prefix_parts else ""

            if filename_prefix:
                json_filename = f"{filename_prefix}-{fn}-chunk-{chunk_num:04d}.json"
            else:
                json_filename = f"{fn}-chunk-{chunk_num:04d}.json"

            payload = {
                "id": base64.urlsafe_b64encode(f"{self.blob_path}-{chunk_num}".encode()).decode(),
                "chunk_name": "|".join(path_parts) if path_parts else f"slide_{chunk_num}",
                **chunk_data,
                **md,
                "toc_detected": False
            }

            try:
                blob_path = os.path.join(bp, json_filename).replace("\\", "/")
                print(f"[DEBUG] Uploading to blob: {blob_path}")
                blob_client = self.dest_container_client.get_blob_client(blob_path)
                blob_client.upload_blob(
                    json.dumps(payload, indent=2), overwrite=True
                )
                print(f"  [UPLOADED] {json_filename}")
            except Exception as e:
                print(f"  [ERROR] Upload failed for {json_filename}: {e}")

    def _extract_metadata(self):
        blob_path = self.blob_path
        parts = blob_path.split("/")
        if len(parts) >= 4:
            return {"Client": parts[0], "Project": parts[1], "Module": parts[2], "Source": parts[3], "File": "/".join(parts[4:])}
        elif len(parts) == 3:
            return {"Client": parts[0], "Project": parts[1], "Module": "N/A", "Source": "N/A", "File": parts[2]}
        else:
            return {"Client": "N/A", "Project": "N/A", "Module": "N/A", "Source": "N/A", "File": parts[0] if parts else "N/A"}

    def _log_step(self, step_name, message, level="INFO"):
        """Logging utility"""
        timestamp = time.strftime("%H:%M:%S")
        print(f"[{timestamp}][{level}][{step_name}] {message}")

    # ── [PATCH] Native PPTX adaptation ──────────────────────────────────────
    def _process_native_pptx_slides(self, native_slides, slide_notes, openai_available):
        """
        Adapt slides produced by robust_extractor_patched.extract_pptx_native()
        to the chunk format expected by _upload_chunks().

        The native extractor preserves italic, color, source-order images, slide
        durations from notes, and avoids OCR pollution. This helper just wires
        its output into the existing pipeline.

        Args:
            native_slides: list of dicts from extract_pptx_native()
            slide_notes: dict {page_num: notes_text} from _extract_slide_notes()
            openai_available: whether to run LLM structure analysis

        Returns: list of chunks (same format as the DI path)
        """
        slides_data = {}
        for s in native_slides:
            page_num = s.get("page_number") or s.get("source_page_range")
            try:
                page_num = int(str(page_num).split('-')[0])
            except (ValueError, TypeError):
                continue
            slides_data[page_num] = {
                "title": s.get("topic", ""),
                "subtopics": s.get("subtopics", "N/A"),
                "content": s.get("content", ""),
                "images": s.get("images", []),
                "tables": s.get("tables", []),
                "_slide_duration": s.get("_slide_duration"),
                "_slide_notes": s.get("_slide_notes", ""),
            }
        total_slides = len(slides_data)
        self._log_step("NATIVE_PROCESSING", f"Adapted {total_slides} native slides")

        # LLM structure analysis (same path as DI)
        llm_structure = {}
        if openai_available and total_slides:
            self._log_step("LLM_ANALYSIS", f"Starting LLM structure analysis for {total_slides} slides...")
            try:
                llm_structure = self._analyze_structure_with_llm(total_slides, slides_data) or {}
            except Exception as e:
                self._log_step("LLM_ANALYSIS", f"LLM analysis failed: {e}", "WARNING")

        # Build chunks
        chunks = []
        for i in range(1, total_slides + 1):
            slide = slides_data.get(i, {})
            content = slide.get('content', '')
            if i in llm_structure:
                final_chapter = llm_structure[i].get("chapter", "Introduction")
                final_topic = llm_structure[i].get("topic", slide.get('title', f"Slide {i}"))
                final_subtopics = llm_structure[i].get("subtopics", slide.get('subtopics', "N/A"))
            else:
                final_chapter = "Introduction"
                final_topic = slide.get('title', f"Slide {i}")
                final_subtopics = slide.get('subtopics', "N/A")

            # Append slide notes (kept separate from main content for the CO,
            # but still passed through so downstream consumers can use them)
            notes = slide_notes.get(i) or slide.get('_slide_notes', '')
            if notes:
                content_with_notes = content + f"\n\n--- SLIDE NOTES ---\n{notes}"
            else:
                content_with_notes = content

            chunks.append({
                "page_number": i,
                "chapter": final_chapter,
                "topic": final_topic,
                "subtopics": final_subtopics,
                "content": content_with_notes,
                "images": slide.get('images', []),
                "tables": [],
                "source_page_range": str(i),
                "duration": slide.get('_slide_duration', ''),
            })

        # Upload chunks via the standard pipeline
        self._upload_chunks(chunks)
        self._log_step("NATIVE_PROCESSING", f"Native PPTX path complete. {len(chunks)} chunks uploaded.")
        return chunks
    # ────────────────────────────────────────────────────────────────────────


# --- 1. STANDARD PRESENTATION AGENT ---

class PresentationAgent(BaseAgent):
    def __init__(self, config, local_doc_path, blob_path, toc_data=None, toc_page_range=None):
        # Pass parameters to the updated BaseAgent __init__
        super().__init__(config, local_doc_path, blob_path, toc_data, toc_page_range)
        self._log_step("INIT", "Standard Presentation Agent initialized.")
        self._log_step("INIT", f"OpenAI client available: {self.openai_client is not None}")

    #----old code -----------------------------------------------------------

    # def _extract_slide_notes(self):
    #     self._log_step("SLIDE_NOTES", "Extracting slide notes (checking for hidden slides)...")
    #     notes_by_slide = defaultdict(str)
    #     try:
    #         if self.local_doc_path.lower().endswith('.pptx'):
    #             prs = Presentation(self.local_doc_path)
    #             visible_slide_index = 0
    #             for slide in prs.slides:
    #                 if slide.element.get("show") == "0": continue
    #                 visible_slide_index += 1
    #                 if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
    #                     note_text = slide.notes_slide.notes_text_frame.text
    #                     if note_text.strip(): 
    #                         notes_by_slide[visible_slide_index] = note_text
    #             self._log_step("SLIDE_NOTES", f"Found notes for {len(notes_by_slide)} slides")
    #     except Exception as e:
    #         self._log_step("SLIDE_NOTES", f"Slide note extraction failed: {e}", "WARNING")
    #     return notes_by_slide

    # -- NEW CODE:  which exctract notes with better formatting and bullet point detection, also checking for hidden slides

    def _extract_slide_notes(self):
        self._log_step("SLIDE_NOTES", "Extracting slide notes (checking for hidden slides)...")
        notes_by_slide = defaultdict(str)
        try:
            if self.local_doc_path.lower().endswith('.pptx'):
                prs = Presentation(self.local_doc_path)
                visible_slide_index = 0
                for slide in prs.slides:
                    if slide.element.get("show") == "0":
                        continue
                    visible_slide_index += 1
                    if slide.has_notes_slide:
                        notes_tf = slide.notes_slide.notes_text_frame
                        lines = []
                        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                        for para in notes_tf.paragraphs:
                            text = para.text.strip()
                            if not text:
                                continue
                            indent = "  " * para.level
                            # Check XML for bullet character or auto-number
                            pPr = para._p.pPr
                            has_bullet = False
                            if pPr is not None:
                                if (pPr.find(f'{{{ns}}}buChar') is not None or
                                        pPr.find(f'{{{ns}}}buAutoNum') is not None):
                                    has_bullet = True
                            if has_bullet:
                                lines.append(f"{indent}• {text}")
                            else:
                                lines.append(f"{indent}{text}")
                        note_text = "\n".join(lines)
                        if note_text.strip():
                            notes_by_slide[visible_slide_index] = note_text
                self._log_step("SLIDE_NOTES", f"Found notes for {len(notes_by_slide)} slides")
        except Exception as e:
            self._log_step("SLIDE_NOTES", f"Slide note extraction failed: {e}", "WARNING")
        return notes_by_slide
    

    def _clean_title(self, text):
        if not text: return "N/A"
        text = re.sub(r'^\s*(section|chapter|module|part|day|week|unit)\s*\d+[\s:\-]*', '', text, flags=re.IGNORECASE)
        text = re.sub(r'^\s*\d+[\.\)]\s+', '', text)
        text = re.sub(r'^\s*[:\-|–]\s*', '', text)
        return re.sub(r'\s+', ' ', text).strip()

    def _detect_repeating_elements(self, di_result, threshold=8):
        self._log_step("REPEATING_ELEMENTS", f"Detecting recurring headers/footers (Threshold: >{threshold} slides)...")
        text_occurrence_by_page = defaultdict(set)
        
        if di_result.paragraphs:
            for para in di_result.paragraphs:
                clean_text = para.content.strip()
                if not clean_text: continue
                if para.bounding_regions:
                    page_num = para.bounding_regions[0].page_number
                    text_occurrence_by_page[clean_text].add(page_num)

        ignored_texts = set()
        for text, pages in text_occurrence_by_page.items():
            if len(pages) > threshold:
                ignored_texts.add(text)
        
        self._log_step("REPEATING_ELEMENTS", f"Found {len(ignored_texts)} unique recurring text patterns to exclude.")
        return ignored_texts

    def _call_openai_api(self, system_prompt, user_prompt):
        """Make the actual OpenAI API call - USING SAME PATTERN AS HIERARCHICALCHUNKINGAGENT"""
        self._log_step("OPENAI_API", "Making OpenAI API call...")
        try:
            response = self.openai_client.chat.completions.create(
                model=self.openai_deployment,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.0,
                response_format={"type": "json_object"}
            )
            result = response.choices[0].message.content
            self._log_step("OPENAI_API", f"OpenAI API call successful ({len(result)} chars)")
            return result
        except Exception as e:
            self._log_step("OPENAI_API", f"OpenAI API call failed: {e}", "ERROR")
            raise

    def _analyze_structure_with_llm(self, total_slides, slides_data):
        self._log_step("LLM_ANALYSIS", f"Analyzing presentation structure using LLM (Total Slides: {total_slides})...")
        
        if not self.openai_client:
            self._log_step("LLM_ANALYSIS", "No Azure OpenAI client configured. LLM structure analysis will be skipped.", "ERROR")
            self._log_step("LLM_ANALYSIS", "self.openai_client is None", "DEBUG")
            return {}

        self._log_step("LLM_ANALYSIS", "OpenAI client is available, proceeding with LLM analysis")

        # 1. Prepare Content for LLM
        context_lines = []
        for i in range(1, total_slides + 1):
            slide = slides_data.get(i, {})
            title = slide.get('title', 'No Title')
            content_snippet = slide.get('content', '')[:300].replace("\n", " ")
            context_lines.append(f"Slide {i}: [Raw Title: {title}] | Content Snippet: {content_snippet}")
        
        full_context = "\n".join(context_lines)
        self._log_step("LLM_ANALYSIS", f"Prepared context for LLM ({len(full_context)} chars)")

        # 2. UPDATED PROMPT: Emphasize Chronological Flow
        system_prompt = (
            "You are an AI assistant specialized in structuring presentation content. "
            "Your task is three-fold: "
            "1. Assign a 'chapter' to each slide based on the chronological flow. "
            "   IMPORTANT: Do not group non-adjacent slides into one chapter if there is a different topic in between. "
            "   Example: If 'Quiz' appears at Slide 10 and Slide 50, but Slide 20 is 'Biology', name them 'Quiz 1' and 'Quiz 2', or just 'Knowledge Checks' but ensure the JSON order is correct. "
            "2. Determine a concise, descriptive 'topic' for EACH slide. "
            "3. Generate a concise 'subtopics' string summarizing key bullets. "
            "Return a strict JSON object containing a list of objects, one for every slide."
        )

        user_prompt = f"""
        Here is the slide list (Total: {total_slides} slides):
        
        {full_context}

        Output strict JSON in this format:
        {{
            "structure": [
                {{ 
                    "slide_number": 1, 
                    "chapter": "Introduction", 
                    "topic": "Welcome & Overview",
                    "subtopics": "1. Introductions\\n2. Agenda" 
                }},
                {{ 
                    "slide_number": 2, 
                    "chapter": "Module 1", 
                    "topic": "Core Concepts",
                    "subtopics": "1. Definition of Terms\\n2. Key Examples" 
                }}
            ]
        }}
        
        Ensure every slide from 1 to {total_slides} is included in order.
        """

        try:
            # Use ThreadPoolExecutor with timeout for Function App environment
            self._log_step("LLM_ANALYSIS", "Starting LLM analysis with timeout...")
            with ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(self._call_openai_api, system_prompt, user_prompt)
                try:
                    # Set timeout appropriate for Function App (Azure Functions have max timeout limits)
                    result_text = future.result(timeout=180)  # 3 minutes timeout
                    self._log_step("LLM_ANALYSIS", f"LLM response received ({len(result_text)} chars)")
                    
                    data = json.loads(result_text)
                    structure_list = data.get("structure", [])
                    
                    self._log_step("LLM_ANALYSIS", f"Parsed {len(structure_list)} slide structures from LLM")
                    
                    structure_map = {}
                    for item in structure_list:
                        try:
                            s_num = int(item.get("slide_number"))
                            structure_map[s_num] = {
                                "chapter": item.get("chapter", "Unassigned"),
                                "topic": item.get("topic", "General"),
                                "subtopics": item.get("subtopics", "N/A")
                            }
                        except ValueError as ve:
                            self._log_step("LLM_ANALYSIS", f"Invalid slide number in LLM response: {item.get('slide_number')} - {ve}", "WARNING")
                            continue

                    if len(structure_map) != total_slides:
                        self._log_step("LLM_ANALYSIS", f"LLM returned {len(structure_map)} slides, expected {total_slides}", "WARNING")
                    
                    return structure_map

                except FutureTimeoutError:
                    self._log_step("LLM_ANALYSIS", "LLM API call timed out after 3 minutes", "ERROR")
                    return {}
                except Exception as e:
                    self._log_step("LLM_ANALYSIS", f"Error in LLM analysis thread: {e}", "ERROR")
                    return {}

        except json.JSONDecodeError as je:
            self._log_step("LLM_ANALYSIS", f"Failed to parse LLM JSON response: {je}", "ERROR")
            if 'result_text' in locals():
                self._log_step("LLM_ANALYSIS", f"Response text: {result_text[:500]}...", "DEBUG")
            return {}
        except Exception as e:
            self._log_step("LLM_ANALYSIS", f"LLM Structure detection failed: {e}", "ERROR")
            self._log_step("LLM_ANALYSIS", f"Full traceback: {traceback.format_exc()}", "DEBUG")
            return {}

    def _process_di_result(self, di_result, extracted_images=None):
        """
        Process DI results into per-slide content with inline image markers.

        Args:
            di_result: the Document Intelligence result
            extracted_images: optional dict {page_num: [image_dict, ...]} from
                _extract_images_and_screenshots(). When provided, inline
                [IMAGE: id] markers reference the deduplicated image IDs
                (e.g., p1_img_1, p1_img_2 — not raw DI figure indices).
                Without this, markers fall back to figure-index numbering
                which may not match the deduplicated image set.
        """
        slides_data = {}
        self._log_step("DI_PROCESSING", "Processing DI content with header exclusion...")

        if extracted_images is None:
            extracted_images = {}

        ignored_texts = self._detect_repeating_elements(di_result, threshold=8)

        figures_by_page = defaultdict(list)
        if di_result.figures:
            for figure in di_result.figures:
                if figure.bounding_regions:
                    page_num = figure.bounding_regions[0].page_number
                    figures_by_page[page_num].append(figure)

        # [PATCH] Pre-compute figure rects per page for OCR-overlap filtering.
        # When DI runs OCR on images, it returns the OCR text as paragraphs
        # whose bounding regions lie INSIDE the figure region. We drop those
        # paragraphs so OCR text doesn't pollute the slide content.
        # Fixes: "No content from image should be taken. It should always be
        #         an image" (Comm p23) and similar.
        figure_rects_by_page = defaultdict(list)
        if di_result.figures:
            for figure in di_result.figures:
                if not figure.bounding_regions:
                    continue
                page_num = figure.bounding_regions[0].page_number
                poly = figure.bounding_regions[0].polygon
                if not poly or len(poly) < 8:
                    continue
                xs = poly[0::2]
                ys = poly[1::2]
                figure_rects_by_page[page_num].append({
                    'x0': min(xs), 'y0': min(ys),
                    'x1': max(xs), 'y1': max(ys),
                })

        def _para_inside_figure(para, page_num):
            """Return True if the paragraph's bbox is inside any figure on this page."""
            if not para.bounding_regions:
                return False
            poly = para.bounding_regions[0].polygon
            if not poly or len(poly) < 8:
                return False
            xs = poly[0::2]
            ys = poly[1::2]
            px0, py0 = min(xs), min(ys)
            px1, py1 = max(xs), max(ys)
            p_area = max((px1 - px0) * (py1 - py0), 0.0001)
            for r in figure_rects_by_page.get(page_num, []):
                ix0 = max(px0, r['x0'])
                iy0 = max(py0, r['y0'])
                ix1 = min(px1, r['x1'])
                iy1 = min(py1, r['y1'])
                if ix1 > ix0 and iy1 > iy0:
                    overlap = (ix1 - ix0) * (iy1 - iy0)
                    if overlap / p_area > 0.7:
                        return True
            return False

        tables_by_page = defaultdict(list)
        if di_result.tables:
            for table in di_result.tables:
                if table.bounding_regions:
                    page_num = table.bounding_regions[0].page_number
                    tables_by_page[page_num].append(table)

        for page_idx, page in enumerate(di_result.pages):
            page_num = page_idx + 1

            valid_paragraphs = []

            if di_result.paragraphs:
                # [PATCH] Sliding-window paragraph dedup catches DI re-emitting
                # the same text from overlapping bounding regions, which is
                # the root cause of "Duplicate text" / "Word repeated twice"
                # reviewer comments.
                seen_recent = []
                for paragraph in di_result.paragraphs:
                    if (paragraph.bounding_regions and
                        paragraph.bounding_regions[0].page_number == page_num):

                        para_text = paragraph.content.strip()

                        if para_text in ignored_texts:
                            continue
                        # Skip header and footer roles explicitly
                        if paragraph.role in ['pageHeader', 'pageFooter', 'pageNumber']:
                            continue

                        # [PATCH] Skip paragraphs that are OCR from inside images
                        if _para_inside_figure(paragraph, page_num):
                            continue

                        if not para_text:
                            continue

                        # [PATCH] Sliding-window dedup against recent paragraphs
                        if para_text in seen_recent:
                            continue
                        seen_recent.append(para_text)
                        if len(seen_recent) > 8:
                            seen_recent.pop(0)

                        valid_paragraphs.append(paragraph)

#--OLD CODE FOR TABLE FORMATION ------------------------------------------------------

            # selected_title = None
            # subheadings = [] 
            # clean_body_texts = []
            
            # for p in valid_paragraphs:
            #     txt = p.content.strip()
            #     role = p.role if p.role else ""
                
            #     if role == 'title':
            #         if selected_title is None: selected_title = txt
            #         else: subheadings.append(txt)
            #     elif role in ['sectionHeading', 'heading']: subheadings.append(txt)
            #     else: clean_body_texts.append(txt)

            # if not selected_title:
            #     if subheadings: selected_title = subheadings.pop(0)
            #     elif clean_body_texts:
            #         for para in clean_body_texts:
            #             if 5 <= len(para) <= 100 and not para.startswith('[IMAGE'):
            #                 selected_title = para
            #                 break
            #         else: selected_title = clean_body_texts[0] if clean_body_texts else f"Slide {page_num}"
            #     else: selected_title = f"Slide {page_num}"

            # cleaned_title = self._clean_title(selected_title)

            # formatted_subtopics = "N/A"
            # if subheadings:
            #     formatted_subtopics = "\n".join([f"{i+1}) {self._clean_title(heading)}" for i, heading in enumerate(subheadings)])
            
            # text_content = "\n\n".join([p.content for p in valid_paragraphs])
            # final_content_parts = [text_content]


            # NEW CODE FOR TABLE AND IMAGE PROCESSING ------------------------------------------------------

            # ── Title extraction ───────────────────────────────────────────
            selected_title       = None
            remaining_paragraphs = []

            for p in valid_paragraphs:
                role = p.role or ""
                if role == 'title' and selected_title is None:
                    selected_title = p.content.strip()
                else:
                    remaining_paragraphs.append(p)

            if not selected_title:
                for i, p in enumerate(remaining_paragraphs):
                    txt  = p.content.strip()
                    role = p.role or ""
                    if role in ('sectionHeading', 'heading'):
                        selected_title       = txt
                        remaining_paragraphs = remaining_paragraphs[i + 1:]
                        break
                    elif 5 <= len(txt) <= 100 and not txt.startswith('[IMAGE'):
                        selected_title       = txt
                        remaining_paragraphs = remaining_paragraphs[i + 1:]
                        break
                if not selected_title:
                    selected_title = f"Slide {page_num}"

            cleaned_title = self._clean_title(selected_title)

            # ── Find leftmost x of list items for nesting depth ────────────
            list_x_min = None
            for p in remaining_paragraphs:
                if p.role == 'listItem' and p.bounding_regions:
                    poly = p.bounding_regions[0].polygon
                    if poly:
                        x = min(poly[0::2])
                        if list_x_min is None or x < list_x_min:
                            list_x_min = x
            if list_x_min is None:
                list_x_min = 0.0

            # ── Build formatted content ────────────────────────────────────
            content_lines = []
            subheadings   = []
            in_list       = False

            for p in remaining_paragraphs:
                txt  = p.content.strip()
                role = p.role or ""

                if role in ('sectionHeading', 'heading'):
                    if in_list:
                        content_lines.append("")
                        in_list = False
                    subheadings.append(txt)
                    content_lines.append(f"\n**{txt}**")

                elif role == 'listItem':
                    in_list = True
                    # Nesting level from bounding box x-offset (~0.3" per level)
                    level = 0
                    if p.bounding_regions:
                        poly = p.bounding_regions[0].polygon
                        if poly:
                            x_left = min(poly[0::2])
                            level  = max(0, round((x_left - list_x_min) / 0.3))
                            level  = min(level, 4)
                    # Strip any leading dash/bullet chars DI may embed
                    clean_txt = re.sub(r'^[-*•·◦▸▹►▻»]\s*', '', txt)
                    clean_txt = re.sub(r'^\d+[.)]\s*', '', clean_txt)
                    clean_txt = re.sub(r':(?:un)?selected:\s*', '', clean_txt).strip()

                # OLD CODE FOR BULLET POINTS ------------------------------------------------------    
                    # indent = "  " * level
                    # bullet = "•" if level == 0 else ("◦" if level == 1 else "▪")
                    # content_lines.append(f"{indent}{bullet} {clean_txt}")

                # NEW CODE FOR BULLET POINTS ------------------------------------------------------
                    indent = "  " * level   # keep indentation
                    content_lines.append(f"{indent}- {clean_txt}")   # use - instead of bullet chars

                else:
                    # Regular body text — DI markdown preserves **bold** and *italic*
                    if in_list:
                        content_lines.append("")
                        in_list = False
                    if txt:
                        content_lines.append(txt)
                        content_lines.append("")

            # ── Subtopics from section headings ────────────────────────────
            formatted_subtopics = "N/A"
            if subheadings:
                formatted_subtopics = "\n".join(
                    f"{i + 1}) {self._clean_title(h)}" for i, h in enumerate(subheadings)
                )

            # ── Assemble final content ──────────────────────────────────────
            content_str = "\n".join(content_lines)
            content_str = re.sub(r'\n{3,}', '\n\n', content_str).strip()

            # Append table markdown
            for table in tables_by_page.get(page_num, []):
                table_md = self._format_di_table(table)
                if table_md:
                    content_str += f"\n\n{table_md}"

            # ── [PATCH] Inline image positioning ───────────────────────────
            # OLD: All [IMAGE: id] markers were appended at the END of content,
            #      after all text and tables. Renderer faithfully placed them
            #      at the end. Reviewer comments: "Icons placed above text
            #      instead of correct position", "Images not placed properly
            #      as per the source".
            #
            # NEW: For each EXTRACTED image (after dedup), compute its
            #      y-position on the page from the matching DI figure and
            #      insert the [IMAGE: id] marker between the paragraphs whose
            #      y-positions bracket the figure. This preserves source order
            #      AND uses the same IDs that the renderer will resolve.
            #
            # If no images survive dedup, we just emit the content_str as-is.
            page_image_records = extracted_images.get(page_num, []) if extracted_images else []
            if not page_image_records:
                final_content = content_str
            else:
                # Build (y_top, marker_id) tuples for SURVIVING images.
                # Each extracted image carries _y_top from the dedup pass.
                # Fall back to evenly-distributed positions if y is missing.
                figure_markers = []
                for img_rec in page_image_records:
                    y_top = img_rec.get('_y_top')
                    if y_top is None:
                        # Without position info, just append at end
                        y_top = 1.0
                    figure_markers.append((y_top, f"[IMAGE: {img_rec['id']}]"))

                # Build (y_top, text_block) tuples for paragraphs in this page
                para_blocks = []
                for p in valid_paragraphs:
                    if not p.bounding_regions:
                        continue
                    poly = p.bounding_regions[0].polygon
                    if not poly or len(poly) < 8:
                        continue
                    para_blocks.append((min(poly[1::2]), p.content.strip()))

                if not para_blocks:
                    # No paragraph y-positions to interleave with, fall back
                    # to text-then-images
                    extras = "\n\n".join(m for _, m in figure_markers)
                    final_content = (content_str + ("\n\n" + extras if extras else "")).strip()
                else:
                    # Interleave figures by y-position into the rendered text.
                    # Strategy: rebuild content_str using the same paragraph
                    # ordering as content_lines, but insert image markers at
                    # the y-position breakpoints.
                    #
                    # Simpler approach that matches the existing content_lines
                    # structure: split content_str on paragraphs (line-level),
                    # then walk both sequences in y-order and emit markers
                    # between paragraphs whose y straddles the figure.
                    figure_markers.sort(key=lambda x: x[0])
                    para_blocks.sort(key=lambda x: x[0])

                    # Build a list of content lines with images interleaved.
                    # We use the original content_lines (in document order)
                    # because they already have the proper bullet/heading
                    # formatting. We map each line back to a paragraph y if
                    # possible, then insert markers at the right boundaries.
                    para_y_by_text = {p.content.strip(): min(p.bounding_regions[0].polygon[1::2])
                                      for p in valid_paragraphs
                                      if p.bounding_regions and p.bounding_regions[0].polygon}

                    out_lines = []
                    fig_idx = 0
                    last_y = -1.0
                    for line in content_lines:
                        # Strip formatting prefixes to find the paragraph y
                        clean = re.sub(r'^[\s\-•◦▪]*\**', '', line).strip().rstrip('*')
                        # Try exact and substring match
                        line_y = para_y_by_text.get(clean)
                        if line_y is None:
                            for k, v in para_y_by_text.items():
                                if clean and (clean in k or k in clean):
                                    line_y = v
                                    break

                        # Emit any figure markers whose y is between last_y and line_y
                        if line_y is not None:
                            while fig_idx < len(figure_markers) and figure_markers[fig_idx][0] <= line_y:
                                out_lines.append("")
                                out_lines.append(figure_markers[fig_idx][1])
                                out_lines.append("")
                                fig_idx += 1
                            last_y = line_y
                        out_lines.append(line)

                    # Emit any remaining figure markers (below the last paragraph)
                    while fig_idx < len(figure_markers):
                        out_lines.append("")
                        out_lines.append(figure_markers[fig_idx][1])
                        out_lines.append("")
                        fig_idx += 1

                    final_content = "\n".join(out_lines)
                    final_content = re.sub(r'\n{3,}', '\n\n', final_content).strip()

                    # Re-append table markdown that was added to content_str
                    # but isn't in content_lines (tables come from a different
                    # source loop). Detect if content_str had table data and
                    # append it at the end (before any trailing images).
                    table_md_appended = []
                    for table in tables_by_page.get(page_num, []):
                        tmd = self._format_di_table(table)
                        if tmd:
                            table_md_appended.append(tmd)
                    if table_md_appended:
                        final_content += "\n\n" + "\n\n".join(table_md_appended)

            slides_data[page_num] = {
                "title":     cleaned_title,
                "subtopics": formatted_subtopics,
                "content":   final_content.strip(),
                "tables":    tables_by_page.get(page_num, [])
            }


        self._log_step("DI_PROCESSING", f"Processed {len(slides_data)} slides from DI")
        return slides_data

    # STEP 6b — Format a DI table as a markdown pipe-table
    # ------------------------------------------------------------------
    def _format_di_table(self, table):
        """Converts a DI table object into a markdown pipe-table string."""
        if not table.cells:
            return ""

# OLD CODE FOR TABLE FORMATION ------------------------------------------------------
        # rows = defaultdict(dict)
        # for cell in table.cells:
        #     rows[cell.row_index][cell.column_index] = cell.content.strip()

# NEW CODE FOR TABLE FORMATION ------------------------------------------------------

        rows = defaultdict(dict)
        for cell in table.cells:
            content = cell.content.strip()
            row_span = getattr(cell, 'row_span', 1) or 1
            col_span = getattr(cell, 'column_span', 1) or 1
            for dr in range(row_span):
                for dc in range(col_span):
                    r = cell.row_index + dr
                    c = cell.column_index + dc
                    # Only top-left cell of merge gets the content, rest are empty
                    rows[r][c] = content if (dr == 0 and dc == 0) else ""

        if not rows:
            return ""

        max_row = max(rows.keys())
        max_col = max(max(cols.keys()) for cols in rows.values())

        lines = []
        for r in range(max_row + 1):
            row_cells = [rows.get(r, {}).get(c, "") for c in range(max_col + 1)]
            lines.append("| " + " | ".join(row_cells) + " |")
            if r == 0:
                lines.append("|" + "|".join([" --- "] * (max_col + 1)) + "|")

        return "\n".join(lines)
    
# ----------------------------------------- till here added new code -------------------------------------------


    def _extract_images_and_screenshots(self, pdf_path, di_result):
        # [PATCH] Major rework to fix multiple reviewer comments:
        #
        #   "A few icons are repeated"               -> dedupe by xref per page
        #   "Repeated icons"                          -> dedupe by xref per page
        #   "Icons or images are duplicated"          -> dedupe by xref per page
        #   "Icons not correct as in source"          -> tighter overlap match
        #   "Image missing"                           -> softer header/footer filter
        #   "First image not required"                -> position-repeat detection
        #
        # Strategy:
        #   1. For each DI figure, find the BEST-overlapping PDF image xref
        #      (instead of first-found with >50% overlap). Pick the xref with
        #      maximum (intersection / figure_area).
        #   2. Deduplicate by xref per page so the same image isn't extracted
        #      twice from overlapping figures.
        #   3. Soften the header/footer filter from 10% to 5%, and only filter
        #      when an image at that position appears on >= 50% of pages
        #      (i.e., it's actually a recurring header/footer, not just a
        #      legitimate icon that happens to be near the top of one slide).
        images_data = defaultdict(list)
        screenshots_data = {}
        try:
            doc = fitz.open(pdf_path)

            # ── Pass 1: collect candidate image positions per page ──
            # Used by the recurring-header detector below.
            page_images_summary = {}  # page_num -> list of (top_y_norm, bottom_y_norm, xref)
            for page_idx in range(len(doc)):
                page = doc.load_page(page_idx)
                page_h = page.rect.height or 1
                summary = []
                try:
                    for img_info in page.get_images(full=True):
                        try:
                            bb = page.get_image_bbox(img_info)
                            summary.append((bb.y0 / page_h, bb.y1 / page_h, img_info[0]))
                        except Exception:
                            pass
                except Exception:
                    pass
                page_images_summary[page_idx + 1] = summary

            # ── Detect recurring-position images (real headers/footers) ──
            # An image is considered a recurring header/footer if it appears at
            # the same vertical position (within ~3%) on >= 50% of pages.
            n_pages = len(doc)
            recurring_positions = set()  # set of rounded y_top values
            position_counts = defaultdict(int)
            for page_num, summary in page_images_summary.items():
                for top, bot, xref in summary:
                    # Only consider top 15% / bottom 15% as candidates
                    if top < 0.15 or bot > 0.85:
                        # Round to nearest 3% bucket
                        bucket = round(top * 33) / 33
                        position_counts[bucket] += 1
            recurring_threshold = max(2, n_pages // 2)
            recurring_positions = {b for b, c in position_counts.items() if c >= recurring_threshold}
            if recurring_positions:
                self._log_step("IMAGE_EXTRACTION",
                    f"Detected {len(recurring_positions)} recurring header/footer position(s)")

            # ── Pass 2: extract figures, with dedup and best-overlap matching ──
            if di_result.figures:
                for fig in di_result.figures:
                    if not fig.bounding_regions:
                        continue
                    page_num = fig.bounding_regions[0].page_number
                    poly = fig.bounding_regions[0].polygon

                    di_page = di_result.pages[page_num - 1]
                    page_height = di_page.height if di_page.height else 1
                    fig_top_norm = min(poly[1::2]) / page_height
                    fig_bottom_norm = max(poly[1::2]) / page_height

                    # ── Soft header/footer filter ──
                    # Only skip if the image is in the top/bottom 5% (was 10%)
                    # AND its position is in the recurring set.
                    is_top = fig_top_norm < 0.05
                    is_bottom = fig_bottom_norm > 0.95
                    if is_top or is_bottom:
                        bucket = round(fig_top_norm * 33) / 33
                        if bucket in recurring_positions:
                            continue
                        # Otherwise, image is in margin but unique to this page — keep it.

                    rect = fitz.Rect(
                        min(poly[0::2]) * 72, min(poly[1::2]) * 72,
                        max(poly[0::2]) * 72, max(poly[1::2]) * 72
                    )
                    rect_area = rect.get_area() or 1

                    try:
                        page = doc.load_page(page_num - 1)
                        # ── Best-overlap matching ──
                        # Iterate ALL images and pick the one with maximum
                        # overlap with this figure's bounding rect, instead of
                        # first-found with >50% overlap.
                        best_xref = None
                        best_overlap = 0.0
                        for img_info in page.get_images(full=True):
                            try:
                                img_bbox = page.get_image_bbox(img_info)
                                overlap = (rect & img_bbox).get_area() / rect_area
                                if overlap > best_overlap:
                                    best_overlap = overlap
                                    best_xref = img_info[0]
                            except Exception:
                                continue

                        # Require >= 30% overlap to count as a match (was >50%
                        # but using best-found is stricter than any threshold)
                        if best_xref is None or best_overlap < 0.30:
                            continue

                        # ── Per-page xref dedup ──
                        # Skip if we've already extracted this xref on this page.
                        already_seen = {
                            img.get("_xref") for img in images_data[page_num]
                        }
                        if best_xref in already_seen:
                            continue

                        base_img = doc.extract_image(best_xref)
                        if base_img:
                            img_count = len(images_data[page_num]) + 1
                            images_data[page_num].append({
                                "id": f"p{page_num}_img_{img_count}",
                                "format": base_img["ext"],
                                "data": base64.b64encode(base_img["image"]).decode('utf-8'),
                                "_xref": best_xref,        # for dedup
                                "_y_top": fig_top_norm,    # for ordering later
                                "_y_bottom": fig_bottom_norm,
                            })
                    except Exception:
                        pass

            for i in range(len(doc)):
                page_num = i + 1
                pix = doc.load_page(i).get_pixmap(dpi=150)
                screenshots_data[page_num] = {
                    "id": f"p{page_num}_slide_view",
                    "format": "png",
                    "data": base64.b64encode(pix.tobytes("png")).decode("utf-8")
                }
            doc.close()
            self._log_step("IMAGE_EXTRACTION",
                f"Extracted {sum(len(v) for v in images_data.values())} images "
                f"and {len(screenshots_data)} screenshots "
                f"(after dedup and recurring-header filter)")
        except Exception as e:
            self._log_step("IMAGE_EXTRACTION", f"Image extraction error: {e}", "ERROR")
        return images_data, screenshots_data

    def _test_openai_connection(self):
        """Test OpenAI connection before processing"""
        if not self.openai_client:
            return False
            
        try:
            self._log_step("OPENAI_TEST", "Testing OpenAI connectivity with quick call...")
            test_response = self.openai_client.chat.completions.create(
                model=self.openai_deployment,
                messages=[{"role": "user", "content": "Reply with 'OK'"}],
                max_tokens=5,
                timeout=10
            )
            result = test_response.choices[0].message.content
            self._log_step("OPENAI_TEST", f"OpenAI test successful: {result}")
            return True
        except Exception as e:
            self._log_step("OPENAI_TEST", f"OpenAI connectivity test failed: {e}", "ERROR")
            return False

    def process(self):
        self._log_step("PROCESS", "Starting Standard Presentation Processing...")

        # Test OpenAI connection first
        openai_available = self._test_openai_connection()
        if not openai_available:
            self._log_step("PROCESS", "OpenAI not available. Will proceed with fallback structure.", "WARNING")

        # Step 1: Extract slide notes
        slide_notes = self._extract_slide_notes()

        # ── [PATCH] Try native PPTX extractor first ───────────────────────
        # The native extractor (robust_extractor_patched.extract_pptx_native)
        # uses python-pptx directly and preserves source structure perfectly.
        # It is dramatically better than DI for:
        #   - Italic + color extraction
        #   - Source-order image positioning
        #   - Per-slide image extraction (no OCR pollution)
        #   - Slide duration parsing from notes
        #
        # We only fall back to the DI-based path if the native path returns
        # empty (e.g., password-protected file) or raises an exception.
        if self.local_doc_path.lower().endswith(('.pptx', '.ppt')):
            try:
                from robust_extractor_patched import extract_pptx_native
                self._log_step("PROCESS", "Trying native PPTX extractor first...")
                native_slides = extract_pptx_native(self.local_doc_path)
                if native_slides:
                    self._log_step("PROCESS",
                        f"Native extraction succeeded ({len(native_slides)} slides). "
                        f"Skipping DI pipeline.")
                    return self._process_native_pptx_slides(native_slides, slide_notes, openai_available)
                else:
                    self._log_step("PROCESS",
                        "Native extraction returned empty. Falling back to DI pipeline.",
                        "WARNING")
            except Exception as e:
                self._log_step("PROCESS",
                    f"Native extraction failed: {e}. Falling back to DI pipeline.",
                    "WARNING")
        # ──────────────────────────────────────────────────────────────────

        # Step 2: Convert to PDF if needed
        pdf_path = self._convert_to_pdf()
        self._log_step("PROCESS", f"Using PDF: {pdf_path}")

        # Step 3: Azure Document Intelligence Analysis
        self._log_step("AZURE_DI", "Analyzing document with Azure Document Intelligence...")
        try:
            with open(pdf_path, "rb") as f:
                poller = self.document_intelligence_client.begin_analyze_document(
                    "prebuilt-layout", AnalyzeDocumentRequest(bytes_source=f.read()),
                    output_content_format=DocumentContentFormat.MARKDOWN)
            di_result = poller.result()
            self._log_step("AZURE_DI", f"DI analysis complete. Pages: {len(di_result.pages) if di_result.pages else 0}")
        except Exception as e:
            self._log_step("AZURE_DI", f"Document Intelligence analysis failed: {e}", "ERROR")
            raise

        # [PATCH] Step 3b: Extract images BEFORE processing DI so the DI
        # processor can use the deduplicated image set when computing inline
        # marker positions. Previously _process_di_result computed markers
        # from raw DI figures (p1_img_1, p1_img_2, ...) but those didn't
        # match the deduplicated extracted_images list, leaving some markers
        # unresolved.
        extracted_images, screenshots = self._extract_images_and_screenshots(pdf_path, di_result)

        # Step 4: Process DI results, passing in deduplicated images
        slides_data = self._process_di_result(di_result, extracted_images=extracted_images)
        total_slides = len(slides_data)
        
        # Step 5: Get Chapters AND Topics from LLM (only if OpenAI is available)
        llm_structure = {}
        if openai_available:
            self._log_step("LLM_ANALYSIS", f"Starting LLM structure analysis for {total_slides} slides...")
            llm_structure = self._analyze_structure_with_llm(total_slides, slides_data)
            
            if not llm_structure:
                self._log_step("LLM_ANALYSIS", "LLM structure analysis failed or returned empty. Using fallback structure.", "WARNING")
            else:
                self._log_step("LLM_ANALYSIS", f"LLM analysis successful. Got structure for {len(llm_structure)} slides")
        else:
            self._log_step("LLM_ANALYSIS", "Skipping LLM analysis due to OpenAI unavailability", "INFO")
        
        # (extracted_images and screenshots already computed above in Step 3b)

        # Step 7: Create chunks
        chunks = []
        for i in range(1, total_slides + 1):
            slide = slides_data.get(i, {})
            content = slide.get('content', '')
            
            # Determine Final Chapter and Topic
            if i in llm_structure:
                final_chapter = llm_structure[i].get("chapter", "Introduction")
                final_topic = llm_structure[i].get("topic", slide.get('title', f"Slide {i}"))
                final_subtopics = llm_structure[i].get("subtopics", slide.get('subtopics', "N/A"))
                self._log_step("CHUNK_CREATION", f"Slide {i}: Using LLM Data -> Chapter: {final_chapter}, Topic: {final_topic}")
            else:
                final_chapter = "Introduction"
                final_topic = slide.get('title', f"Slide {i}")
                final_subtopics = slide.get('subtopics', "N/A")
                self._log_step("CHUNK_CREATION", f"Slide {i}: Using Raw Data -> Topic: {final_topic}")
            
            # Append Notes
            if i in slide_notes:
                content += f"\n\n--- SLIDE NOTES ---\n{slide_notes[i]}"

            # Append Screenshot ref
            # [PATCH] Skip full-slide screenshots in CO output.
            # Reviewer global comment (Param Sharma, multiple rows): "Please remove
            # snips of all slides as shown in all rows wherever seen. The slides
            # snips are not required." Screenshots are still kept in the images
            # list for debugging/reference purposes, but they are NOT embedded
            # into the content field where they'd render as visual duplicates
            # of the source slide.
            ss = screenshots.get(i)
            imgs = extracted_images.get(i, [])
            if ss:
                ss['id'] = f"p{i}_full_slide"
                ss['_is_full_slide_screenshot'] = True  # Mark for downstream filtering
                # Do NOT append [IMAGE: p{i}_full_slide] to content — this was
                # causing the screenshot to be embedded in the CO output.
                # imgs.append(ss)  # Also do NOT add to imgs list since the
                #                    renderer would embed it as a visible image.
                # The screenshot remains available in the `screenshots` dict
                # for any downstream process that needs it (e.g., Layer 2 vision).

            # [PATCH] Strip private dedup-tracking keys (_xref, _y_top, _y_bottom)
            # before uploading. They were used internally for positioning and
            # dedup but should not pollute the search index.
            clean_imgs = [
                {k: v for k, v in img.items() if not k.startswith('_')}
                for img in imgs
            ]

            chunks.append({
                "page_number": i,
                "chapter": final_chapter,
                "topic": final_topic,
                "subtopics": final_subtopics,
                "content": content,
                "images": clean_imgs,
                "tables": [],
                "source_page_range": str(i)
            })

        # Print Structure Summary
        print("\n" + "="*80)
        print(f"{'DETECTED PRESENTATION STRUCTURE':^80}")
        print("="*80)

        current_chapter_name = None
        current_chapter_slides = []

        # Helper to print the buffer
        def print_buffer(c_name, c_slides):
            if c_name and c_slides:
                start = c_slides[0]['page_number']
                end = c_slides[-1]['page_number']
                print(f"\n[CHAPTER] {c_name.upper()} (Page Range: {start}-{end})")
                print("-" * 40)
                for s in c_slides:
                    print(f"  • Slide {s['page_number']}: {s['topic']}")

        for chunk in chunks:
            chapter = chunk['chapter']
            
            # If chapter changes, print the previous buffer and start new
            if chapter != current_chapter_name:
                print_buffer(current_chapter_name, current_chapter_slides)
                current_chapter_name = chapter
                current_chapter_slides = [chunk]
            else:
                current_chapter_slides.append(chunk)

        # Print the final buffer
        print_buffer(current_chapter_name, current_chapter_slides)
        
        print("\n" + "="*80 + "\n")

        # Step 8: Upload chunks
        self._upload_chunks(chunks)
        
        # Step 9: Cleanup temporary PDF if created
        if pdf_path != self.local_doc_path:
            try:
                os.remove(pdf_path)
                self._log_step("CLEANUP", f"Cleaned up temporary PDF: {pdf_path}")
            except Exception as e:
                self._log_step("CLEANUP", f"Failed to cleanup temporary PDF: {e}", "WARNING")
        
        self._log_step("PROCESS", "Presentation processing completed successfully!")


# --- 2. NARRATION PPT AGENT (PDF with Slide+Text) ---

# --- 2. NARRATION PPT AGENT (PDF with Slide+Text) ---

class NarrationPptAgent(BaseAgent):
    def __init__(self, config, local_doc_path, blob_path, toc_data=None, toc_page_range=None):
        super().__init__(config, local_doc_path, blob_path, toc_data, toc_page_range)
        self._log_step("INIT", "Narration PDF Agent initialized.")
        
        # Test OpenAI connection for structure analysis
        self.openai_available = self._test_openai_connection()

    def _test_openai_connection(self):
        """Test OpenAI connection for structure analysis"""
        if not self.openai_client:
            return False
            
        try:
            self._log_step("OPENAI_TEST", "Testing OpenAI connectivity for narration PDF...")
            test_response = self.openai_client.chat.completions.create(
                model=self.openai_deployment,
                messages=[{"role": "user", "content": "Reply with 'OK'"}],
                max_tokens=5,
                timeout=10
            )
            result = test_response.choices[0].message.content
            self._log_step("OPENAI_TEST", f"OpenAI test successful: {result}")
            return True
        except Exception as e:
            self._log_step("OPENAI_TEST", f"OpenAI connectivity test failed: {e}", "ERROR")
            return False

    def _extract_slide_image(self, pdf_path, page_num, di_figures):
        """Extract the slide image from a page"""
        if not di_figures: 
            return None
            
        try:
            doc = fitz.open(pdf_path)
            page = doc.load_page(page_num - 1)
            
            # Try to find the largest figure that looks like a slide
            best_fig = None
            best_area = 0
            
            for fig in di_figures:
                if fig.bounding_regions:
                    poly = fig.bounding_regions[0].polygon
                    width = poly[2] - poly[0]
                    height = poly[5] - poly[1]
                    
                    # Check if it's slide-sized (typically wide aspect ratio)
                    if width > 5.0 and width/height > 1.2:  # Wider than tall
                        area = width * height
                        if area > best_area:
                            best_area = area
                            best_fig = fig
            
            if not best_fig:
                # If no slide-like figure, use the page screenshot
                pix = page.get_pixmap(dpi=150)
                data = base64.b64encode(pix.tobytes("png")).decode("utf-8")
                doc.close()
                return {
                    "id": f"slide_{page_num}",
                    "format": "png", 
                    "data": data,
                    "description": f"Page {page_num} screenshot"
                }
            
            # Extract the specific figure region
            poly = best_fig.bounding_regions[0].polygon
            rect = fitz.Rect(
                poly[0] * 72, 
                poly[1] * 72, 
                poly[4] * 72, 
                poly[5] * 72
            )
            
            # Add padding
            rect.x0 -= 5
            rect.y0 -= 5
            rect.x1 += 5
            rect.y1 += 5
            
            pix = page.get_pixmap(clip=rect, dpi=150)
            data = base64.b64encode(pix.tobytes("png")).decode("utf-8")
            doc.close()
            
            return {
                "id": f"slide_{page_num}",
                "format": "png", 
                "data": data,
                "description": f"Slide image from page {page_num}"
            }
        except Exception as e:
            self._log_step("IMAGE_EXTRACTION", f"Failed to extract slide image: {e}", "WARNING")
            return None

    def _is_slide_image_page(self, page_num, di_figures):
        """Determine if this page contains a slide image (not just any figure)"""
        if not di_figures:
            return False
            
        # Look for figures that are likely slides (wide and large)
        for fig in di_figures:
            if fig.bounding_regions:
                poly = fig.bounding_regions[0].polygon
                width = poly[2] - poly[0]
                height = poly[5] - poly[1]
                
                # Slides are typically wide (aspect ratio > 1.2) and reasonably sized
                if width > 5.0 and width/height > 1.2:
                    return True
        
        return False

    def _extract_text_from_page(self, page_spans, di_result):
        """Extract text from a page"""
        page_text = ""
        if page_spans:
            for span in page_spans:
                page_text += di_result.content[span.offset : span.offset + span.length]
        return page_text

    def _analyze_slide_structure_with_llm(self, slide_data):
        """Use LLM to analyze slide structure and extract chapters/topics"""
        if not self.openai_available or not slide_data:
            return {}
        
        self._log_step("LLM_ANALYSIS", f"Analyzing {len(slide_data)} slides with LLM...")
        
        # Prepare slide information for LLM
        context_lines = []
        for page_num, data in slide_data.items():
            text_preview = data.get('text', '')[:200].replace('\n', ' ')
            context_lines.append(f"Page {page_num}: {text_preview}")
        
        full_context = "\n".join(context_lines)
        
        system_prompt = """You are an expert at analyzing presentation content from narration PDFs.
        Your task:
        1. Identify logical chapters/sections in the presentation
        2. Extract meaningful topic titles for each slide
        3. Maintain chronological order
        
        Return JSON in this format:
        {
            "structure": [
                {
                    "slide_number": 1,
                    "chapter": "Introduction",
                    "topic": "Welcome & Course Overview"
                }
            ]
        }"""
        
        user_prompt = f"""Analyze this presentation content from a PDF with slide images and narration text.
        The content is extracted from pages that contain slide images.
        
        Content preview:
        {full_context}
        
        Provide chapter and topic assignments for each slide page. Focus on the main content, not navigation or repeated elements.
        """
        
        try:
            response = self.openai_client.chat.completions.create(
                model=self.openai_deployment,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.0,
                response_format={"type": "json_object"}
            )
            
            result = response.choices[0].message.content
            data = json.loads(result)
            structure_map = {}
            
            for item in data.get("structure", []):
                slide_num = item.get("slide_number")
                if slide_num:
                    structure_map[int(slide_num)] = {
                        "chapter": item.get("chapter", "Unassigned"),
                        "topic": item.get("topic", "General")
                    }
            
            self._log_step("LLM_ANALYSIS", f"LLM analysis complete: {len(structure_map)} slides analyzed")
            return structure_map
            
        except Exception as e:
            self._log_step("LLM_ANALYSIS", f"LLM analysis failed: {e}", "WARNING")
            return {}

    def _extract_title_from_slide_text(self, slide_text):
        """Extract a meaningful title from slide text"""
        if not slide_text:
            return "Untitled"
        
        lines = [line.strip() for line in slide_text.split('\n') if line.strip()]
        
        # Look for potential titles (short lines, not too long)
        for line in lines[:5]:  # Check first 5 lines
            if 5 <= len(line) <= 100:
                # Clean common prefixes
                clean_line = re.sub(r'^\s*(slide|page)\s*\d+[\.:]\s*', '', line, flags=re.IGNORECASE)
                clean_line = re.sub(r'^\s*\d+[\.\)]\s*', '', clean_line)
                clean_line = re.sub(r'^\s*[:\-|•]\s*', '', clean_line)
                
                if clean_line.strip():
                    return clean_line.strip()[:80]
        
        # If no good title found, use first non-empty line
        for line in lines:
            if line:
                return line[:50]
        
        return "Untitled"

    def process(self):
        self._log_step("PROCESS", "Starting Narration-style PDF Processing...")
        
        # Step 1: Convert to PDF if needed
        pdf_path = self._convert_to_pdf()
        self._log_step("PROCESS", f"Using PDF: {pdf_path}")
        
        # Step 2: Analyze with Document Intelligence
        with open(pdf_path, "rb") as f:
            poller = self.document_intelligence_client.begin_analyze_document(
                "prebuilt-layout", 
                AnalyzeDocumentRequest(bytes_source=f.read()),
                output_content_format=DocumentContentFormat.MARKDOWN
            )
            result = poller.result()
        
        self._log_step("DI", f"Analysis complete. Pages: {len(result.pages)}")
        
        # Step 3: Organize figures by page
        figures_by_page = defaultdict(list)
        if result.figures:
            for fig in result.figures:
                if fig.bounding_regions:
                    page_num = fig.bounding_regions[0].page_number
                    figures_by_page[page_num].append(fig)
        
        # Step 4: Identify slide pages and text pages
        slide_pages = {}
        text_pages_by_slide = defaultdict(list)
        
        # Get page spans for text extraction
        page_spans = {}
        if result.pages:
            for page in result.pages:
                page_spans[page.page_number] = page.spans
        
        # First pass: Identify all slide pages
        slide_page_numbers = []
        for page_num in sorted(page_spans.keys()):
            if self._is_slide_image_page(page_num, figures_by_page.get(page_num, [])):
                slide_page_numbers.append(page_num)
                # Extract slide text
                slide_text = self._extract_text_from_page(page_spans[page_num], result)
                slide_pages[page_num] = {
                    "text": slide_text,
                    "title": self._extract_title_from_slide_text(slide_text)
                }
        
        self._log_step("SLIDE_DETECTION", f"Found {len(slide_pages)} slide pages: {slide_page_numbers}")
        
        # Step 5: Group text pages with their preceding slide
        current_slide = None
        for page_num in sorted(page_spans.keys()):
            if page_num in slide_pages:
                current_slide = page_num
            elif current_slide:
                # This is a text page following a slide
                text_pages_by_slide[current_slide].append(page_num)
        
        # Log grouping information
        for slide_page, text_pages in text_pages_by_slide.items():
            self._log_step("GROUPING", f"Slide page {slide_page} has {len(text_pages)} text pages: {text_pages}")
        
        # Step 6: Get LLM structure analysis for slides
        llm_structure = self._analyze_slide_structure_with_llm(slide_pages)
        
        # Step 7: Create chunks (one per slide + its narration text)
        chunks = []
        slide_counter = 1
        
        for slide_page_num in sorted(slide_pages.keys()):
            # Get slide image
            slide_img = self._extract_slide_image(
                pdf_path, 
                slide_page_num, 
                figures_by_page.get(slide_page_num, [])
            )
            
            # Determine chapter and topic
            if slide_page_num in llm_structure:
                chapter = llm_structure[slide_page_num].get("chapter", "Introduction")
                topic = llm_structure[slide_page_num].get("topic", slide_pages[slide_page_num]["title"])
            else:
                # Fallback: Use detected title
                chapter = "Day 0: 5"  # Default from your output
                topic = slide_pages[slide_page_num]["title"]
            
            # Get slide text
            slide_text = slide_pages[slide_page_num]["text"]
            
            # Get narration text from following pages
            narration_text = ""
            text_page_nums = text_pages_by_slide.get(slide_page_num, [])
            for text_page_num in text_page_nums:
                page_text = self._extract_text_from_page(page_spans[text_page_num], result)
                narration_text += f"\n\n--- Page {text_page_num} ---\n{page_text}"
            
            # Combine slide text and narration
            full_content = f"# {topic}\n\n"
            if slide_img:
                full_content += f"[IMAGE: {slide_img['id']}]\n\n"
            
            full_content += f"## Slide Content\n{slide_text}"
            
            if narration_text.strip():
                full_content += f"\n\n## Narration\n{narration_text}"
            
            # Create chunk
            chunk = {
                "slide_number": slide_counter,
                "page_number": slide_page_num,
                "chapter": chapter,
                "topic": topic,
                "content": full_content,
                "images": [slide_img] if slide_img else [],
                "tables": [],
                "source_page_range": f"{slide_page_num}" + (f"-{text_page_nums[-1]}" if text_page_nums else ""),
                "text_pages": text_page_nums
            }
            
            chunks.append(chunk)
            slide_counter += 1
        
        self._log_step("PROCESS", f"Created {len(chunks)} chunks from {len(slide_pages)} slides")
        
        # Step 8: Log structure like PresentationAgent
        print("\n" + "="*80)
        print(f"{'NARRATION PDF STRUCTURE':^80}")
        print("="*80)
        
        current_chapter_name = None
        current_chapter_slides = []
        
        def print_buffer(c_name, c_slides):
            if c_name and c_slides:
                start = c_slides[0]['slide_number']
                end = c_slides[-1]['slide_number']
                print(f"\n[CHAPTER] {c_name.upper()} (Slides: {start}-{end})")
                print("-" * 40)
                for s in c_slides:
                    text_pages = s.get('text_pages', [])
                    page_info = f"Page {s['page_number']}"
                    if text_pages:
                        page_info += f" + {len(text_pages)} text pages"
                    print(f"  • Slide {s['slide_number']}: {s['topic']} ({page_info})")
        
        for chunk in chunks:
            chapter = chunk['chapter']
            
            if chapter != current_chapter_name:
                print_buffer(current_chapter_name, current_chapter_slides)
                current_chapter_name = chapter
                current_chapter_slides = [chunk]
            else:
                current_chapter_slides.append(chunk)
        
        print_buffer(current_chapter_name, current_chapter_slides)
        print("="*80 + "\n")
        
        # Step 9: Upload chunks
        self._upload_chunks(chunks)
        
        # Step 10: Cleanup
        if pdf_path != self.local_doc_path:
            try:
                os.remove(pdf_path)
                self._log_step("CLEANUP", f"Cleaned up temporary PDF: {pdf_path}")
            except Exception as e:
                self._log_step("CLEANUP", f"Failed to cleanup PDF: {e}", "WARNING")
        
        self._log_step("PROCESS", "Narration PDF processing completed!")


# --- MAIN ENTRY POINT ---

# --- MAIN ENTRY POINT ---

if __name__ == "__main__":
    # Run environment diagnostics
    check_function_app_environment()
    
    parser = argparse.ArgumentParser(description="Fixed Presentation Processing Agent")
    parser.add_argument("--blob-path", required=True, help="Full path in the blob storage")
    parser.add_argument("--local-path", required=True, help="Local path to the presentation file")
    parser.add_argument("--config", required=True, help="JSON string with configuration details")
    args = parser.parse_args()

    try:
        print("\n" + "="*80)
        print("STARTING PRESENTATION PROCESSING")
        print("="*80)
        
        # Parse config
        print("[DEBUG] Parsing configuration...")
        config = json.loads(args.config)
        
        # Debug: Show what config keys are available
        print(f"[DEBUG] Config has {len(config)} keys")
        print(f"[DEBUG] Config keys: {list(config.keys())}")
        
        # Set defaults if not provided (for backward compatibility)
        if "AZURE_DI_ENDPOINT" not in config:
            print("[WARN] AZURE_DI_ENDPOINT not in config, using default")
            config["AZURE_DI_ENDPOINT"] = "https://aptara-di.cognitiveservices.azure.com/"
        
        if "AZURE_DI_KEY" not in config:
            print("[WARN] AZURE_DI_KEY not in config, using default")
            config["AZURE_DI_KEY"] = os.getenv("AZURE_DI_KEY", "")  # SECURITY: removed hardcoded fallback
        
        # Check OpenAI configuration - IMPORTANT: Use same pattern as HierarchicalChunkingAgent
        if "AZURE_OPENAI_ENDPOINT" not in config or "AZURE_OPENAI_API_KEY" not in config:
            print("[WARN] AZURE_OPENAI_ENDPOINT or AZURE_OPENAI_API_KEY missing from config.")
            print("[WARN] Checking environment variables...")
            
            # Check environment variables (for Function App)
            if os.environ.get("AZURE_OPENAI_ENDPOINT") and os.environ.get("AZURE_OPENAI_API_KEY"):
                print("[INFO] Found OpenAI config in environment variables")
                config["AZURE_OPENAI_ENDPOINT"] = os.environ["AZURE_OPENAI_ENDPOINT"]
                config["AZURE_OPENAI_API_KEY"] = os.environ["AZURE_OPENAI_API_KEY"]
            else:
                print("[WARN] OpenAI config not found in config or environment")
                print("[WARN] LLM Chapter/Topic detection will be skipped.")
        else:
            print("[INFO] OpenAI configuration found in config")
            if not config["AZURE_OPENAI_ENDPOINT"] or not config["AZURE_OPENAI_API_KEY"]:
                print("[WARN] OpenAI endpoint or key is empty in config")
            else:
                print("[INFO] OpenAI configuration is valid")
        
        # Set OpenAI deployment name if not provided
        if "AZURE_OPENAI_DEPLOYMENT_NAME" not in config:
            config["AZURE_OPENAI_DEPLOYMENT_NAME"] = os.environ.get("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o")
        
        # Check required configuration
        required_keys = ["AZURE_DI_ENDPOINT", "AZURE_DI_KEY", "STORAGE_CONNECTION_STRING", "DESTINATION_CONTAINER_NAME"]
        missing_keys = [k for k in required_keys if k not in config]
        
        if missing_keys:
            raise ValueError(f"Missing required config keys: {missing_keys}")
        
        print("[DEBUG] All required config keys are present")
        
        # Verify local file exists
        if not os.path.exists(args.local_path):
            raise FileNotFoundError(f"Local file not found: {args.local_path}")
        
        print(f"[DEBUG] Local file exists: {args.local_path}")
        print(f"[DEBUG] Blob path: {args.blob_path}")
        
        # Determine file type and instantiate appropriate agent
        file_extension = os.path.splitext(args.local_path)[1].lower()
        
        if file_extension in ['.pptx', '.ppt']:
            print("[INFO] Using PresentationAgent for PowerPoint file")
            agent = PresentationAgent(
                config=config, 
                local_doc_path=args.local_path, 
                blob_path=args.blob_path,
                toc_data={},
                toc_page_range=None
            )
        elif file_extension in ['.pdf']:
            # Check if it's a narration-style PDF by examining the content
            print("[INFO] Using NarrationPptAgent for PDF file")
            agent = NarrationPptAgent(
                config=config,
                local_doc_path=args.local_path,
                blob_path=args.blob_path,
                toc_data={},
                toc_page_range=None
            )
        else:
            # For other file types (doc, docx, etc.), use PresentationAgent
            # They will be converted to PDF first
            print(f"[INFO] Using PresentationAgent for {file_extension} file")
            agent = PresentationAgent(
                config=config,
                local_doc_path=args.local_path,
                blob_path=args.blob_path,
                toc_data={},
                toc_page_range=None
            )
        
        # Process the document
        agent.process()

        print("\n" + "="*80)
        print("[SUCCESS] Agent completed successfully!")
        print("="*80)

    except json.JSONDecodeError as je:
        print(f"\n[FATAL ERROR] Failed to parse config JSON: {je}")
        print(f"[DEBUG] Config string: {args.config[:200]}...")
        exit(1)
    except FileNotFoundError as fnfe:
        print(f"\n[FATAL ERROR] File not found: {fnfe}")
        exit(1)
    except ValueError as ve:
        print(f"\n[FATAL ERROR] Configuration error: {ve}")
        exit(1)
    except Exception as e:
        print(f"\n[FATAL ERROR] Processing failed: {e}")
        print("\n[DEBUG] Full traceback:")
        traceback.print_exc()
        exit(1)

